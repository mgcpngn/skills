
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Image Paths from previous generation
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"

def create_ppt_v3():
    prs = Presentation()
    
    BG_COLOR = RGBColor(15, 15, 25) # Deep LeiJun Blue/Black

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_copyright(slide, is_cover=False):
        # Mandatory Copyright Injection
        text = "天翼云湖北分公司 作者：王理" if is_cover else "天翼云湖北分公司 王理"
        
        # Position: Bottom Center/Right
        left = Inches(0)
        top = Inches(7.2) # Very bottom
        width = Inches(10)
        height = Inches(0.3)
        
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100) # Subtle Grey
        p.alignment = PP_ALIGN.CENTER

    def add_slide(title_text, content_lines=[], image_path=None, is_title_slide=False, is_split_layout=False):
        layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(layout)
        set_background(slide)

        # Title Logic
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p = title_tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if is_title_slide:
            p.alignment = PP_ALIGN.CENTER
        
        CONTENT_TOP = Inches(2.5) 
        
        if is_title_slide:
            # --- Cover Slide Layout Fix: Image Constraint ---
            # Max available height for image = (Subtitle Top - Title Bottom)
            # Let's say Subtitle starts at 5.5 inches
            SUBTITLE_TOP = Inches(5.8)
            IMG_MAX_HEIGHT = Inches(3.5) # 5.8 - 2.3 = 3.5 roughly
            
            if image_path and os.path.exists(image_path):
                # We specify height to ensure it fits, let width be proportional
                left = Inches(1)
                top = Inches(2.1)
                # By specifying height, we prevent vertical overflow into subtitle
                slide.shapes.add_picture(image_path, left, top, height=IMG_MAX_HEIGHT)
                
            # Subtitles below image
            sub_box = slide.shapes.add_textbox(Inches(1), SUBTITLE_TOP, Inches(8), Inches(1))
            tf = sub_box.text_frame
            for line in content_lines:
                p = tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(200, 200, 200)
            
            # Copyright
            add_copyright(slide, is_cover=True)

        elif is_split_layout:
            # --- Split Layout Safe Zones ---
            text_box = slide.shapes.add_textbox(Inches(0.5), CONTENT_TOP, Inches(4.5), Inches(4.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for line in content_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(14)
                if line.startswith("##"):
                    p.text = line.replace("##", "").strip()
                    p.font.size = Pt(36)
                    p.font.color.rgb = RGBColor(255, 215, 0)
                    p.font.bold = True
                elif line.startswith("#"):
                    p.text = line.replace("#", "").strip()
                    p.font.size = Pt(28)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                else:
                    p.text = line
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(220, 220, 220)
            
            if image_path and os.path.exists(image_path):
                img_left = Inches(5.5)
                img_width = Inches(4.0)
                # Constraint height here too just in case
                slide.shapes.add_picture(image_path, img_left, CONTENT_TOP, width=img_width)

        else:
            # Standard List Layout
            text_box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, Inches(8), Inches(4.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for line in content_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(20)
                if line.startswith("##"):
                    p.text = line.replace("##", "").strip()
                    p.font.size = Pt(54)
                    p.font.color.rgb = RGBColor(255, 215, 0)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                elif line.startswith("#"):
                    p.text = line.replace("#", "").strip()
                    p.font.size = Pt(32)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                else:
                    p.text = "• " + line
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(220, 220, 220)

    # --- Content ---

    # Slide 1: Title
    add_slide("内存的战争\nThe War of Memory", ["2026 内存架构深度解析", "NVIDIA Rubin vs. DeepSeek Engram"], IMG_TITLE, is_title_slide=True)

    # Slide 2: Agenda
    add_slide("汇报大纲", [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件的暴力美学",
        "3. DeepSeek Engram：软件的魔法",
        "4. 巅峰对决：路线之争",
        "5. 未来展望：Agentic AI"
    ])

    # Slide 3: Intro
    add_slide("2026：范式转移", [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ])

    # Slide 4: The Walls
    add_slide("两堵高墙", [
        "# Memory Wall (内存墙)",
        "权重访问速度限制了推理速度",
        "# Context Wall (上下文墙)",
        "序列长度限制了思考深度",
        "我们需要：更快的读写，更大的容量"
    ])

    # Slide 5: The Two Players
    add_slide("两大阵营", [
        "# NVIDIA (硬件派)",
        "集成电路，专有互联",
        "# DeepSeek (软件派)",
        "算法优化，软件定义内存"
    ])

    # Slide 6: NVIDIA Rubin
    add_slide("NVIDIA Rubin 架构", [
        "Jensen Huang @ CES 2026",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN, is_split_layout=True)

    # Slide 7: Rubin GPU Specs
    add_slide("Rubin GPU：性能怪兽", [
        "## 3360 亿",
        "晶体管数量 (Transistors)",
        "## 50 PFLOPS",
        "FP4 推理性能"
    ])

    # Slide 8: HBM4 Revolution
    add_slide("HBM4：速度的革命", [
        "## 22 TB/s",
        "内存带宽 (Bandwidth)",
        "## 288 GB",
        "单卡显存容量"
    ])

    # Slide 9: Vera CPU
    add_slide("Vera CPU：内存扩展", [
        "专为 AI Factory 设计",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "NVLink-C2C 互联 @ 1.8 TB/s"
    ])

    # Slide 11: The Cost
    add_slide("成本与代价", [
        "高性能 = 高投入",
        "专有硬件确立护城河",
        "CapEx (资本支出) 巨大"
    ])

    # Slide 13: DeepSeek Engram
    add_slide("DeepSeek Engram", [
        "软件定义的“外挂内存”",
        "无需专有硬件",
        "利用系统 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM, is_split_layout=True)

    # Slide 14: Core Insight
    add_slide("核心洞察", [
        "LLM 90% 的计算在“重构知识”",
        "为什么不直接“查找”？",
        "# Reconstruct vs. Lookup"
    ])

    # Slide 15: Implementation
    add_slide("实现原理：Engram", [
        "N-gram Hashing (N元语法哈希)",
        "Embedding Table 存放在 DDR",
        "Context-Aware Gating"
    ])

    # Slide 16: Prefetching Strategy
    add_slide("克服 PCIe 瓶颈", [
        "PCIe 比 NVLink 慢得多",
        "解决方案：",
        "## Deterministic Prefetching",
        "异步预取，吞吐量损失 < 3%"
    ])

    # Slide 17: Performance
    add_slide("性能实测：以小博大", [
        "## +5.0",
        "BBH (复杂推理)",
        "## +3.4",
        "MMLU (知识问答)"
    ])

    # Slide 19: Economic Impact
    add_slide("经济学：AI 民主化", [
        "无需 H100/Rubin",
        "消费级显卡 + 大内存",
        "RTX 4090 集群跑 GPT-5"
    ])

    # Slide 20: Strategic Comparison
    add_slide("路线对比：逻辑之争", [
        "# NVIDIA: Cohesion",
        "硬件透明，高性能",
        "# DeepSeek: Retrieval",
        "软件显式，高性价比"
    ], IMG_COMPARE, is_split_layout=True)

    # Slide 24: Conclusion
    add_slide("总结：内存革命", [
        "Rubin 定义了上限",
        "Engram 拉高了下限",
        "## 遗忘的时代结束了"
    ])

    # Slide 25: Closing (Title Slide style for closing)
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_background(slide)
    
    # Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    p = t_box.text_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    p = sub.text_frame.add_paragraph()
    p.text = "The Architecture of Memory"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    add_copyright(slide, is_cover=False)

    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V3.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")

if __name__ == "__main__":
    create_ppt_v3()
