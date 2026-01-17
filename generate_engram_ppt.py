"""
DeepSeek Engram - 20 Pages PPT
突出 Engram 的价值和意义
应用 leijunskill v2.0 全部技巧
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_NGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_ngram_hashing_1768530510183.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_prefetch_strategy_1768530526058.png"

# --- CONSTANTS ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)

# --- FONT SIZES (leijunskill §2.4) ---
TITLE_SIZE = Pt(54)
BIG_NUMBER_SIZE = Pt(120)
SUBTITLE_SIZE = Pt(32)
BODY_SIZE = Pt(26)

# Colors
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(255, 215, 0)
LIGHT_GRAY = RGBColor(220, 220, 220)
TIANYI_BLUE = RGBColor(0, 191, 255)
DEEPSEEK_GREEN = RGBColor(46, 204, 113)


def create_engram_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    def add_signature(slide, is_cover=True):
        box = slide.shapes.add_textbox(
            (SLIDE_WIDTH - Inches(5)) / 2,
            SLIDE_HEIGHT - Inches(1),
            Inches(5), Inches(0.9)
        )
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = TIANYI_BLUE
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(24)
        p2.font.color.rgb = GOLD
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def content_slide(title, bullets, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Content
        content_box = s.shapes.add_textbox(Inches(1), Inches(2.0), SLIDE_WIDTH - Inches(2), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(14)
            if line.startswith("##"):
                para.text = line[2:].strip()
                para.font.size = BIG_NUMBER_SIZE
                para.font.color.rgb = GOLD
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = SUBTITLE_SIZE
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = "• " + line
                para.font.size = BODY_SIZE
                para.font.color.rgb = LIGHT_GRAY
        add_notes(s, notes)
        return s

    def data_slide(title, number, unit, explanation, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1.2))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Big Number
        num_box = s.shapes.add_textbox(Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(2))
        p = num_box.text_frame.add_paragraph()
        p.text = number
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = DEEPSEEK_GREEN
        p.alignment = PP_ALIGN.CENTER
        # Unit
        unit_box = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
        p = unit_box.text_frame.add_paragraph()
        p.text = unit
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Explanation
        exp_box = s.shapes.add_textbox(Inches(1), Inches(5.2), SLIDE_WIDTH - Inches(2), Inches(1))
        p = exp_box.text_frame.add_paragraph()
        p.text = explanation
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_notes(s, notes)
        return s

    def split_slide(title, bullets, img_path, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Left content
        left_box = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(12)
            if line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = Pt(28)
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = line
                para.font.size = Pt(22)
                para.font.color.rgb = LIGHT_GRAY
        # Right image
        if img_path and os.path.exists(img_path):
            s.shapes.add_picture(img_path, Inches(6.5), Inches(2.0), width=Inches(6.3))
        add_notes(s, notes)
        return s

    def closing_slide(title, subtitle, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        sub = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
        p = sub.text_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_signature(s, is_cover=False)
        add_notes(s, notes)
        return s

    # ==================== 20 PAGES: ENGRAM DEEP DIVE ====================

    # --- P1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), SLIDE_WIDTH - Inches(1), Inches(1.8))
    p = box.text_frame.add_paragraph()
    p.text = "DeepSeek Engram\n条件记忆：稀疏性的新维度"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if os.path.exists(IMG_ENGRAM):
        s.shapes.add_picture(IMG_ENGRAM, Inches(4), Inches(2.3), height=Inches(2.8))
    sub = s.shapes.add_textbox(Inches(0), Inches(5.4), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "Conditional Memory via Scalable Lookup | arXiv 2601.07372"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)
    add_notes(s, """朋友们，大家好！

今天，我想给大家介绍一项可能改变AI未来的技术——DeepSeek Engram。

这不是一个普通的技术改进，而是一个全新的范式。它问了一个简单但深刻的问题：

"AI为什么要浪费算力去'回忆'那些早已知道的事实？"

想象一下：你每次回答"北京是中国的首都"，都要从头推导一遍。多累啊！直接查表不行吗？

Engram说：可以。

这就是今天的主角——条件记忆，稀疏性的新维度。""")

    # --- P2: Agenda ---
    content_slide("分享大纲", [
        "1. 问题：AI的失忆与浪费",
        "2. 洞察：重构 vs. 查表",
        "3. 方案：Engram 技术架构",
        "4. 效果：性能与效率双赢",
        "5. 意义：AI的未来范式"
    ], """我们今天的分享分为五个部分。

首先，我们看看当前大语言模型存在什么问题——为什么它们既"失忆"又浪费算力。

然后，我来揭示DeepSeek团队的核心洞察——一个简单却被忽视的真理。

接着，我们深入Engram的技术架构，看看它是如何工作的。

之后，我用数据说话，展示Engram带来的性能提升。

最后，我们一起思考：Engram对AI的未来意味着什么。

准备好了吗？Let's go！""")

    # --- P3: The Problem - AI Amnesia ---
    content_slide("问题：AI的失忆症", [
        "上下文窗口有限",
        "KV Cache 指数膨胀",
        "聊着聊着就忘了",
        "## AI 记不住你说过的话"
    ], """（人设塑造：真诚的奋斗者）

朋友们，你们有没有遇到过这种情况？

你和AI聊了半个小时，聊得很开心。突然你问它："我刚才说的那个想法，你觉得怎么样？"

它愣住了，说："抱歉，你说过什么想法吗？"

我有个同事，用AI写季度报告。写到一半，让AI修改第一段。AI回答："抱歉，我不记得第一段写了什么。"

他当场崩溃。

这就是AI的失忆症。上下文窗口有限，聊得越多，忘得越多。这是一个让所有人头疼的问题。""")

    # --- P4: The Problem - Computational Waste ---
    content_slide("问题：算力的浪费", [
        "90% 计算在'重构'静态知识",
        "北京是中国首都——每次都要重新推导",
        "1+1=2——每次都要重新计算",
        "## 浪费，是 LLM 的原罪"
    ], """更糟糕的是，AI还在浪费大量算力。

（生活化比喻）
想象你每次回答"北京是中国的首都"，都要从头推导一遍。先想亚洲有哪些国家，再想中国的地理，最后得出结论。

这不荒谬吗？这个事实你早就知道了，直接说不行吗？

但现在的大语言模型就是这么干的！研究表明，90%的计算都花在"重构"那些早已知道的静态知识上。

北京是首都、水是H2O、1+1=2……这些事实每次都要重新"推导"。

这是巨大的浪费。DeepSeek发现了这个问题，他们决定解决它。""")

    # --- P5: Core Insight ---
    content_slide("核心洞察", [
        "#重构 (Reconstruction)",
        "用神经网络计算出答案",
        "#查表 (Lookup)",
        "直接从记忆中取出答案",
        "## 别算了，直接查！"
    ], """（峰终定律：这是核心洞察）

DeepSeek的洞察很简单：

现在的LLM在干嘛？重构。用神经网络一层一层计算，最终"推导"出答案。

但这些答案，很多是早就知道的事实啊！

为什么不直接查表呢？

（生活化比喻）
就像考试。天才学霸什么都记在脑子里，考试全靠回忆。但聪明的普通学生带了参考书，脑子记不住的，翻书就行。

结果呢？聪明学生考得可能比学霸还好。因为书上的东西比脑子记得更准确。

这就是Engram的精髓：别算了，直接查！

从"重构"到"查表"，这是范式的转变。""")

    # --- P6: What is Engram ---
    split_slide("什么是 Engram", [
        "#条件记忆模块",
        "Conditional Memory Module",
        "#稀疏性的新维度",
        "与 MoE 互补，而非替代",
        "#O(1) 复杂度查表",
        "无论表多大，查询恒定时间"
    ], IMG_ENGRAM,
    """那么，Engram到底是什么？

Engram是一个"条件记忆模块"，嵌入到Transformer架构中。

它代表了一种新的稀疏性维度。过去，我们有专家混合（MoE）——条件计算。现在，我们有Engram——条件记忆。

（雷氏对比法）
MoE解决的是"用哪些神经元计算"的问题。
Engram解决的是"查哪些记忆"的问题。

两者不是替代关系，而是互补关系。加在一起，效果更强。

关键特性是什么？O(1)复杂度。无论你的知识表有多大——100亿、1000亿参数——查询时间都是恒定的。

这是一个优雅的数学解决方案。""")

    # --- P7: N-gram Hashing ---
    split_slide("技术：N-gram Hashing", [
        "#输入 Token 编码为 Hash ID",
        "确定性、上下文相关",
        "#用 Hash 作为索引",
        "直接定位知识向量",
        "#多头哈希 (Multi-head)",
        "捕捉不同层次的模式"
    ], IMG_NGRAM,
    """具体怎么做的呢？

核心技术叫 N-gram Hashing。

（生活化比喻）
想象你有一本超级大的百科全书，1000亿个词条。现在我问"北京"，你要找到关于北京的内容。

如果从头翻到尾？翻到天荒地老也翻不完。

聪明的办法：给每个词条一个编号，按编号存放。我说"北京"，你算出编号是"12345"，直接翻到第12345页。

N-gram Hashing就是这么干的。把输入的Token序列转成一个确定性的数字，用这个数字去查表。

而且是多头哈希——不同的"头"捕捉不同层次的模式，就像用多个索引来查阅。

快得让人难以置信。""")

    # --- P8: Tokenizer Compression ---
    data_slide("词表压缩", "23%", "有效词汇量减少", "通过规范化标识符实现",
    """Engram还做了一个巧妙的优化——词表压缩。

他们发现，很多Token其实是同一个概念的不同形式。比如"AI"、"A.I."、"artificial intelligence"，其实指的是同一个东西。

通过创建规范化的标识符，Engram把有效词汇量减少了23%！

这意味着什么？

（数字具象化）
原来需要100亿个词条，现在只需要77亿。省了23%的存储空间，查询效率也更高。

这就是工程师思维——不放过任何一个优化点。""")

    # --- P9: Memory Hierarchy ---
    content_slide("多级缓存架构", [
        "#热门词条 → GPU HBM",
        "高频访问，驻留显存",
        "#长尾词条 → 主机 DDR",
        "低频访问，按需加载",
        "## Zipf 定律的妙用"
    ], """Engram还利用了一个自然语言的特性——Zipf定律。

什么是Zipf定律？简单说：少数词语出现频率极高，大多数词语很少出现。

（生活化比喻）
就像你的手机App。微信、抖音每天用，放在首屏。某个三年前下载的App，从来不用，藏在最后一屏。

Engram同样这么做：

热门词条（高频访问）→ 放在GPU显存里，随时取用。
长尾词条（低频访问）→ 放在主机DDR内存里，需要时再加载。

这样，既保证了速度，又节省了昂贵的GPU显存。

这就是工程的智慧——顺应自然规律，而不是逆天而行。""")

    # --- P10: Prefetching ---
    split_slide("确定性预取", [
        "#计算先于需求",
        "提前算出 Hash ID",
        "#异步加载数据",
        "GPU 计算时，内存已在准备",
        "#吞吐损失 < 3%",
        "几乎无感知"
    ], IMG_PREFETCH,
    """但有一个问题：从主机DDR取数据，比从GPU显存慢得多。

DeepSeek怎么解决的？确定性预取。

（生活化比喻）
就像点外卖。你知道自己中午要吃饭，不会等到饿了再点。提前下单，等你饿的时候，外卖已经送到了。

Engram也是这样。因为Hash ID是确定性的，系统可以提前几层就算出来：接下来需要哪些数据。

然后异步加载：GPU在算第1层的时候，内存已经在准备第3层需要的数据。

结果呢？吞吐损失不到3%！

也就是说，用便宜的DDR内存，达到了97%的GPU显存效果。这就是软件的力量。""")

    # --- P11: Key Numbers ---
    data_slide("核心参数", "100B", "卸载到主机内存", "27B 活跃模型 + 100B 知识表",
    """（数字具象化）

100B参数！

这是Engram的知识表规模。100亿参数，全部存放在主机的DDR内存中。

活跃模型只有27B参数——这是实际在GPU上运行的部分。

27B + 100B = 127B。但只有27B在GPU上，其余100B用便宜的DDR承载。

（雷氏对比法：以己之有比彼之无）
传统方案：要跑127B参数模型，需要几百GB显存，价值几十万的GPU。
Engram方案：27B参数的GPU需求 + 几百块钱的DDR内存条。

效果呢？堪比GPT-4级别的知识能力。

这，就是Engram的魔法。""")

    # --- P12: Performance - BBH ---
    data_slide("性能：复杂推理", "+5.0", "BBH 测试提升", "vs 同等 FLOPs 的 MoE 模型",
    """让数据说话。

在BBH（Big Bench Hard）复杂推理测试中，Engram比同等FLOPs的MoE模型高出5分！

（雷氏对比法）
这是什么概念？同样的算力预算，Engram的表现全面碾压。

5分的差距，在AI竞赛中是巨大的。就像高考多考5分，可能从二本变一本。

而且这是在"复杂推理"任务上——需要多步逻辑推导的任务。

为什么Engram能赢？因为它把"静态知识回忆"的工作卸载了。神经网络可以专注于真正的推理，而不是浪费算力在"北京是首都"这种事情上。

专注，才能做到极致。""")

    # --- P13: Performance - Knowledge ---
    data_slide("性能：知识问答", "+3.4", "MMLU 测试提升", "知识检索准确率显著提高",
    """在MMLU知识问答测试中，Engram提升了3.4分。

这个结果完全符合预期。Engram的核心就是增强"记忆"能力，知识问答正是它的主场。

（真诚的脆弱性）
说实话，刚开始看到这个数字的时候，我有点不敢相信。一个27B的小模型，怎么能在知识问答上超过那些动辄上百B参数的巨无霸？

后来我想通了：参数多不等于"记得准"。Engram的外挂知识表，存储的是精心设计的embeddings，精度比神经网络自己"悟"出来的更高。

这就是"显式记忆"的优势。""")

    # --- P14: Performance - Code ---
    data_slide("性能：代码生成", "+3.0", "HumanEval 测试提升", "常见编程模式召回更准确",
    """在HumanEval代码生成测试中，Engram提升了3.0分。

为什么代码任务也能提升？

（洞察分享）
编程有很多"固定模式"。比如遍历数组、读写文件、网络请求……这些都是可以查表的！

传统模型需要"悟"出这些模式，每次都要重新推导。Engram直接把常见的代码模式存在知识表里，需要的时候查出来用。

这就像有经验的程序员vs新手。老手有一整套"代码库"在脑子里，写代码飞快。新手每次都要从头想，慢得多。

Engram让AI也拥有了"老手的经验库"。""")

    # --- P15: U-shaped Scaling Law ---
    content_slide("U形 Scaling Law", [
        "#稀疏分配问题",
        "MoE vs Engram 参数如何分配?",
        "#发现 U形曲线",
        "纯 MoE 和纯 Engram 都不是最优",
        "## 混合分配 > 单一路线"
    ], """DeepSeek还发现了一个重要的规律——U形Scaling Law。

问题是这样的：如果你有固定的"稀疏参数预算"，应该全部给MoE专家？还是全部给Engram记忆？

他们做了大量实验，发现两个极端都不是最优解。最优解是混合分配。

（比喻）
就像一个团队，全是思考者，执行力不足。全是执行者，没人思考。最好的团队是两者兼备。

MoE负责"动态推理"，Engram负责"静态记忆"。各司其职，效果最好。

实验表明：把约20-25%的稀疏预算分配给Engram，其余给MoE，这是最优配比。

这就是科学——不是拍脑袋，而是用数据说话。""")

    # --- P16: Strategic Value - Democratization ---
    content_slide("战略价值：AI民主化", [
        "无需顶级 GPU",
        "RTX 4090 + DDR 即可",
        "打破 NVIDIA 垄断",
        "## 让每个人都用得起 AI"
    ], """（情感共鸣：共同体意识）

朋友们，Engram的意义远不止是技术突破。

它代表了AI的民主化。

过去，要跑大模型，你需要价值几十万的H100。普通人根本玩不起。

有了Engram，你只需要一块RTX 4090，加上几百块钱的DDR内存。就能跑出GPT-4级别的效果。

（真诚的脆弱性）
说实话，这让我很感慨。我们这些普通开发者，也终于能参与AI革命了。

这就像当年个人电脑的普及。以前只有大公司才有计算机，现在人人都有。Engram正在做同样的事——让AI从云端走向每个人的桌面。""")

    # --- P17: Strategic Value - Geopolitics ---
    content_slide("战略价值：自主可控", [
        "面对芯片出口管制",
        "软件弥补硬件差距",
        "DeepSeek V3 仅用 $5.58M 训练",
        "## 以智取胜，More with Less"
    ], """还有一层更深的意义——自主可控。

（真诚的奋斗者）
在当前的国际形势下，很多地方拿不到最新的GPU。这是一个现实的困境。

但Engram告诉我们：硬件不够，软件来凑。

DeepSeek V3的训练成本是多少？558万美元。而GPT-4据说花了1亿美元。

同样的结果，成本只有1/20。这就是"以小博大"的智慧。

（共同体意识）
朋友们，这不仅是技术问题，更是战略问题。Engram证明了一条道路：在资源受限的情况下，依然可以做出世界一流的AI。

环境越艰苦，创新越强劲。这是我们的精神。""")

    # --- P18: Future - Agentic AI ---
    content_slide("未来：Agentic AI", [
        "从 Chatbot 到 Agent",
        "长期记忆成为刚需",
        "Engram = AI的长期记忆引擎",
        "## 记忆，是Agent的灵魂"
    ], """（未来展望）

最后，让我们看向未来。

AI正在从"聊天机器人"进化成"智能代理"。Agent不只回答问题，要替你完成任务：订机票、做报告、写代码、管项目。

这需要什么？长期记忆。

Agent需要记住你的偏好、你之前说过的话、任务的上下文。如果每次都失忆，怎么可能完成复杂任务？

Engram正是解决这个问题的关键技术。它是AI的"长期记忆引擎"。

没有记忆，就没有Agent。Engram，让Agent成为可能。""")

    # --- P19: Summary ---
    content_slide("总结：Engram 的价值", [
        "#技术突破",
        "从'重构'到'查表'，O(1) 复杂度",
        "#性能提升",
        "BBH +5.0，MMLU +3.4，HumanEval +3.0",
        "#战略意义",
        "AI民主化，自主可控"
    ], """（峰终定律：高潮总结）

让我来总结一下Engram的价值。

技术层面：它实现了从"重构"到"查表"的范式转变。O(1)复杂度，用数学的优雅解决了工程的难题。

性能层面：在多个基准测试中全面领先。BBH +5.0，MMLU +3.4，HumanEval +3.0。小模型打败大模型。

战略层面：它代表了AI的民主化和自主可控。让每个人都能用得起AI，让资源受限的环境也能做出一流的创新。

（真诚的奋斗者）
这不是某一个天才的成果，而是DeepSeek整个团队日夜奋战的结晶。他们证明了：智慧，可以弥补资源的不足。

2026年，将被铭记为"条件记忆元年"。从今天起，AI不再需要浪费算力在"回忆"上。

这，就是Engram的意义。""")

    # --- P20: Closing ---
    closing_slide("谢谢", "Conditional Memory via Scalable Lookup",
    """（峰终定律：情怀收尾）

朋友们，这就是今天分享的全部内容。

Engram不仅仅是一个技术模块，它代表了一种新的思维方式：

不要让AI浪费算力在已知的事情上。让它专注于真正需要思考的问题。

这是一个简单的道理，但DeepSeek是第一个把它工程化的团队。

（悬念收尾）
Engram只是开始。DeepSeek说，这将成为下一代稀疏模型的"必备原语"。未来会有更多基于Engram的创新。

让我们一起期待。

感谢大家的聆听！朋友们，我们下次再见！""")

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'DeepSeek_Engram_20pages.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_engram_ppt()
