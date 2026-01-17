
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Image Paths (Using the generated paths)
# Note: These paths are unique to this session. 
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"

def create_ppt():
    prs = Presentation()
    
    # Set Slide Master Background to Dark
    slide_master = prs.slide_master
    background = slide_master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 20) # Almost black

    def add_slide(title_text, content_lines=[], image_path=None, layout_index=1, is_title_slide=False):
        if is_title_slide:
            slide_layout = prs.slide_layouts[0] # Title Slide
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content

        slide = prs.slides.add_slide(slide_layout)
        
        # Access background again to ensure it applies (sometimes master doesn't propagate in simple scripts)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 10, 20) # Deep dark blue/black

        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Image
        if image_path and os.path.exists(image_path):
            # Add image on the right or center depending on content
            if not content_lines:
                # Full screen image or centered large
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                slide.shapes.add_picture(image_path, left, top, width=width)
            else:
                # Content on left, image on right
                left = Inches(5.5)
                top = Inches(2)
                width = Inches(4)
                slide.shapes.add_picture(image_path, left, top, width=width)
                
                # Adjust text box width
                body = slide.shapes.placeholders[1]
                body.width = Inches(4.5)

        # Content
        if content_lines:
            # If standard layout
            if not is_title_slide:
                tf = slide.shapes.placeholders[1].text_frame
                tf.clear() # clear default bullet
                
                for line in content_lines:
                    p = tf.add_paragraph()
                    
                    # Check for "big data" style (starts with *)
                    if line.startswith("##"):
                        p.text = line.replace("##", "").strip()
                        p.font.size = Pt(44)
                        p.font.color.rgb = RGBColor(255, 215, 0) # Gold
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                    elif line.startswith("#"):
                        p.text = line.replace("#", "").strip()
                        p.font.size = Pt(32)
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.font.bold = True
                    else:
                        p.text = line
                        p.font.size = Pt(24)
                        p.font.color.rgb = RGBColor(200, 200, 200)

    # Slide 1: Title
    add_slide("内存的战争\nThe War of Memory", ["2026 内存架构深度解析", "NVIDIA Rubin vs. DeepSeek Engram"], IMG_TITLE)

    # Slide 2: Agenda
    add_slide("汇报大纲", [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件的暴力美学",
        "3. DeepSeek Engram：软件的魔法",
        "4. 巅峰对决：路线之争",
        "5. 未来展望：Agentic AI"
    ])

    # Slide 3: Intro - 2026 Paradigm Shift
    add_slide("2026：范式转移", [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "算力不再是瓶颈",
        "## 记忆才是新的石油"
    ])

    # Slide 4: The Walls
    add_slide("两堵高墙", [
        "# Memory Wall (内存墙)",
        "权重访问速度限制了推理速度",
        "# Context Wall (上下文墙)",
        "序列长度限制了思考深度",
        "我们需要：更快的读写，更大的容量"
    ])

    # Slide 5: The Two Players
    add_slide("两大阵营", [
        "# NVIDIA (硬件派)",
        "Rubin 架构",
        "集成电路，专有互联",
        "# DeepSeek (软件派)",
        "Engram 技术",
        "算法优化，软件定义内存"
    ])

    # Slide 6: NVIDIA Rubin Reveal
    add_slide("NVIDIA Rubin 架构", [
        "Jensen Huang @ CES 2026",
        "Extreme Co-design (极致协同)",
        "将整个数据中心视为一个芯片",
        "六大芯片生态系统"
    ], IMG_RUBIN)

    # Slide 7: Rubin GPU Specs
    add_slide("Rubin GPU：性能怪兽", [
        "## 3360 亿",
        "晶体管数量 (Transistors)",
        "## 50 PFLOPS",
        "FP4 推理性能",
        "台积电 N3P 工艺"
    ])

    # Slide 8: HBM4 Revolution
    add_slide("HBM4：速度的革命", [
        "## 22 TB/s",
        "内存带宽 (Bandwidth)",
        "## 288 GB",
        "单卡显存容量",
        "解决 MoE 专家路由延迟",
        "2.75x 相比 Blackwell"
    ])

    # Slide 9: Vera CPU
    add_slide("Vera CPU：内存扩展", [
        "专为 AI Factory 设计",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "NVLink-C2C 互联 @ 1.8 TB/s",
        "让 CPU 内存像显存一样快"
    ])

    # Slide 10: System Cohesion
    add_slide("系统级协同", [
        "GPU + CPU + DPU",
        "NVL72 机柜系统",
        "## 54 TB",
        "机柜总内存容量",
        "统一内存寻址，透明访问"
    ])

    # Slide 11: The Cost
    add_slide("成本与代价", [
        "高性能 = 高投入",
        "专有硬件确立护城河",
        "适合超大规模云服务商",
        "CapEx (资本支出) 巨大"
    ])

    # Slide 12: Transition to Software
    add_slide("另一种可能？", [
        "如果不买新硬件？",
        "如果是普通显卡？",
        "如果是受限环境？",
        "## Software Defines Everything"
    ])

    # Slide 13: DeepSeek Engram
    add_slide("DeepSeek Engram", [
        "软件定义的“外挂内存”",
        "无需专有硬件",
        "利用系统 DDR 内存 (Side-loading)",
        "打破显存容量限制"
    ], IMG_ENGRAM)

    # Slide 14: Core Insight
    add_slide("核心洞察", [
        "LLM 90% 的计算在“重构知识”",
        "为什么不直接“查找”？",
        "# Reconstruct vs. Lookup",
        "从计算密集转向访存密集"
    ])

    # Slide 15: Implementation
    add_slide("实现原理：Engram", [
        "N-gram Hashing (N元语法哈希)",
        "Embedding Table 存放在 DDR",
        "Context-Aware Gating",
        "直接将知识注入残差流"
    ])

    # Slide 16: Prefetching Strategy
    add_slide("克服 PCIe 瓶颈", [
        "PCIe 比 NVLink 慢得多",
        "解决方案：## Deterministic Prefetching",
        "提前计算 Hash ID",
        "异步预取数据",
        "吞吐量损失 < 3%"
    ])

    # Slide 17: Performance - 1
    add_slide("性能实测：以小博大", [
        "Engram-27B vs MoE同级模型",
        "## +5.0",
        "BBH (复杂推理)",
        "## +3.4",
        "MMLU (知识问答)"
    ])

    # Slide 18: Performance - 2
    add_slide("性能实测：代码与数学", [
        "## +3.0",
        "HumanEval (代码)",
        "## +2.4",
        "MATH (逻辑推导)",
        "更大的知识库，更小的计算量"
    ])

    # Slide 19: Economic Impact
    add_slide("经济学：AI 民主化", [
        "无需 H100/Rubin",
        "消费级显卡 + 大内存",
        "RTX 4090 集群跑 GPT-5 级模型",
        "大幅降低推理成本"
    ])

    # Slide 20: Strategic Comparison
    add_slide("路线对比：逻辑之争", [
        "# NVIDIA: Cohesion (内聚)",
        "硬件透明，高性能，高成本",
        "# DeepSeek: Retrieval (检索)",
        "软件显式，高性价比，高灵活性"
    ], IMG_COMPARE)

    # Slide 21: Memory Type Comparison
    add_slide("针对的记忆类型", [
        "Rubin: 动态 KV Cache (短期记忆)",
        "优化对话上下文",
        "Engram: 静态 Knowledge (长期记忆)",
        "优化世界知识库"
    ])

    # Slide 22: Geopolitics
    add_slide("地缘政治影响", [
        "面对出口管制",
        "Engram 是“突围”方案",
        "用软件弥补硬件差距",
        "More with Less"
    ])

    # Slide 23: Future - Agentic AI
    add_slide("未来：Agentic AI", [
        "从 Chatbot 到 Agent",
        "需要处理百万级 Token",
        "代码库、视频流、完整历史",
        "Context is the new Currency"
    ])

    # Slide 24: Conclusion/Summary
    add_slide("总结：内存革命", [
        "Rubin 定义了上限",
        "Engram 拉高了下限",
        "## 遗忘的时代结束了",
        "AI 拥有了真正的“大脑皮层”"
    ])

    # Slide 25: Closing
    add_slide("谢谢", [
        "The Architecture of Memory",
        "Let Inspiration Never be Forgotten.",
        "## 谢谢大家！"
    ], IMG_TITLE)

    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")

if __name__ == "__main__":
    create_ppt()
