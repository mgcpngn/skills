"""
Memory Architecture Report - V5
Generated following leijunskill standards:
- 30 pages as requested
- Signature only on Cover and Closing pages
- Stylized signature with Tianyi Cloud branding
- Strict SplitImage layout with 1" gap
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Image Paths
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"

# --- CONSTANTS (leijunskill §6.1) ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)

# --- SAFE ZONES (leijunskill §6.2) ---
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(1.5)
CONTENT_TOP = Inches(2.0)
CONTENT_HEIGHT = Inches(4.0)
FOOTER_TOP = Inches(6.5)

# --- SPLIT LAYOUT (leijunskill §6.3) ---
SPLIT_TEXT_LEFT = Inches(0.5)
SPLIT_TEXT_WIDTH = Inches(5.5)
SPLIT_IMG_LEFT = Inches(6.5)
SPLIT_IMG_WIDTH = Inches(6.3)
SPLIT_IMG_MAX_HEIGHT = Inches(4.0)


def create_ppt_v5():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_signature(slide, is_cover=True):
        """Stylized signature - leijunskill §5"""
        sig_width = Inches(5)
        sig_height = Inches(0.9)
        sig_left = (SLIDE_WIDTH - sig_width) / 2
        sig_top = SLIDE_HEIGHT - sig_height - Inches(0.15)

        box = slide.shapes.add_textbox(sig_left, sig_top, sig_width, sig_height)
        tf = box.text_frame

        # Line 1: Company
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(0, 191, 255)  # Sky Blue
        p1.alignment = PP_ALIGN.CENTER

        # Line 2: Author Name (Gold, Bold)
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(255, 215, 0)  # Gold
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def add_title(slide, text, center=False):
        box = slide.shapes.add_textbox(Inches(0.5), TITLE_TOP, SLIDE_WIDTH - Inches(1), TITLE_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if center:
            p.alignment = PP_ALIGN.CENTER

    def add_content_lines(slide, lines):
        box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, SLIDE_WIDTH - Inches(2), CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(18)
            if line.startswith("##"):
                p.text = line[2:].strip()
                p.font.size = Pt(60)
                p.font.color.rgb = RGBColor(255, 215, 0)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(36)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = "• " + line
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(220, 220, 220)

    def add_split_page(slide, lines, img_path):
        # Left Text
        box = slide.shapes.add_textbox(SPLIT_TEXT_LEFT, CONTENT_TOP, SPLIT_TEXT_WIDTH, CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(14)
            if line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(220, 220, 220)
        # Right Image
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, SPLIT_IMG_LEFT, CONTENT_TOP, width=SPLIT_IMG_WIDTH)

    # === PAGE GENERATION (30 pages) ===

    # --- Page 1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "内存的战争\nThe War of Memory", center=True)
    if os.path.exists(IMG_TITLE):
        pic = s.shapes.add_picture(IMG_TITLE, Inches(3.5), Inches(2.3), height=Inches(3.0))
    sub = s.shapes.add_textbox(Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "2026 内存架构深度解析 | NVIDIA Rubin vs. DeepSeek Engram"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)

    # --- Page 2: Agenda ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "汇报大纲")
    add_content_lines(s, [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件霸权",
        "3. DeepSeek Engram：软件革命",
        "4. 巅峰对决：路线之争",
        "5. 未来展望：Agentic AI"
    ])

    # --- Page 3-5: Background (3 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2026：范式转移")
    add_content_lines(s, [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Memory Wall：内存墙")
    add_content_lines(s, [
        "GPU 运算单元高速运转",
        "内存带宽成为数据传输瓶颈",
        "## 算力越强，内存墙越高"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Context Wall：上下文墙")
    add_content_lines(s, [
        "KV Cache 随序列长度指数膨胀",
        "模型无法记住完整对话",
        "## 失忆，是 AI 最大的痛点"
    ])

    # --- Page 6-11: NVIDIA Rubin (6 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVIDIA Rubin 架构")
    add_split_page(s, [
        "#Jensen Huang @ CES 2026",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统",
        "将整个数据中心视为一个芯片"
    ], IMG_RUBIN)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin GPU：性能怪兽")
    add_content_lines(s, [
        "## 3360 亿",
        "晶体管数量 (Transistors)",
        "## 50 PFLOPS",
        "FP4 推理性能"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：速度的革命")
    add_content_lines(s, [
        "## 22 TB/s",
        "内存带宽 (Bandwidth)",
        "相比 Blackwell 提升 2.75x"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin GPU：容量震撼")
    add_content_lines(s, [
        "## 288 GB",
        "单卡 HBM4 显存容量",
        "可容纳 1440 亿参数模型"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Vera CPU：内存扩展")
    add_content_lines(s, [
        "88 颗 Olympus ARM 核心",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "NVLink-C2C 互联 @ 1.8 TB/s"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVL72：超级机柜")
    add_content_lines(s, [
        "72 块 Rubin GPU + 36 块 Vera CPU",
        "## 54 TB",
        "机柜总内存容量",
        "统一内存寻址，透明访问"
    ])

    # --- Page 12-17: DeepSeek Engram (6 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "DeepSeek Engram")
    add_split_page(s, [
        "#软件定义的外挂内存",
        "无需专有硬件",
        "利用系统 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "核心洞察")
    add_content_lines(s, [
        "LLM 90%% 的计算在重构知识",
        "为什么不直接查找?",
        "## Reconstruct vs. Lookup"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "实现原理：N-gram Hashing")
    add_content_lines(s, [
        "将输入 Token 序列编码为确定性 Hash",
        "用 Hash 作为 Key 查询 Embedding",
        "## O(1) 复杂度查表"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Embedding Table：知识的仓库")
    add_content_lines(s, [
        "## 100B 参数",
        "存放在主机 DDR 内存中",
        "让 27B 小模型拥有万亿知识"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "克服 PCIe 瓶颈")
    add_content_lines(s, [
        "PCIe 带宽远低于 NVLink",
        "解决方案：",
        "## Deterministic Prefetching",
        "异步预取，吞吐量损失 < 3%"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "性能实测：以小博大")
    add_content_lines(s, [
        "## +5.0",
        "BBH (复杂推理)",
        "## +3.4",
        "MMLU (知识问答)"
    ])

    # --- Page 18-22: Comparison (5 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "路线对比")
    add_split_page(s, [
        "#NVIDIA: Cohesion (内聚)",
        "硬件透明，高性能，高成本",
        "#DeepSeek: Retrieval (检索)",
        "软件显式，高性价比"
    ], IMG_COMPARE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "目标记忆类型对比")
    add_content_lines(s, [
        "#Rubin: 动态 KV Cache (短期记忆)",
        "优化对话上下文",
        "#Engram: 静态 Knowledge (长期记忆)",
        "优化世界知识库"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "成本与投入")
    add_content_lines(s, [
        "#NVIDIA Rubin",
        "高 CapEx，适合超大规模云厂商",
        "#DeepSeek Engram",
        "低成本，消费级硬件可运行"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "经济学：AI 民主化")
    add_content_lines(s, [
        "无需 H100/Rubin",
        "RTX 4090 集群跑 GPT-5",
        "## 打破 NVIDIA 税"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "地缘政治意义")
    add_content_lines(s, [
        "面对出口管制",
        "Engram 是突围方案",
        "## More with Less"
    ])

    # --- Page 23-26: Future (4 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "未来：Agentic AI")
    add_content_lines(s, [
        "从 Chatbot 到 Agent",
        "多步推理，长期记忆",
        "## Context is the new Currency"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "百万 Token 上下文")
    add_content_lines(s, [
        "Rubin CPX 可摄入完整代码库",
        "将整个仓库作为工作记忆",
        "## 代码理解，一步到位"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "共享记忆池")
    add_content_lines(s, [
        "BlueField-4 DPU 管理 KV Cache",
        "多 Agent 共享，无需重算",
        "## 协作式 AI 成为可能"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "长期一致性")
    add_content_lines(s, [
        "Engram 支持 Episodic/Semantic/Procedural 记忆",
        "在 MemoChat 基准测试中：",
        "## +15 分 (仅用 1% Token)"
    ])

    # --- Page 27-28: Summary (2 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "总结：内存革命")
    add_content_lines(s, [
        "Rubin 定义了性能上限",
        "Engram 拉高了可及下限",
        "## 殊途同归：解决 AI 失忆"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "核心结论")
    add_content_lines(s, [
        "#硬件派: 极致性能，全栈锁定",
        "#软件派: 普惠创新，开源突围",
        "## 遗忘的时代，结束了"
    ])

    # --- Page 29: Full-page Quote ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(1), Inches(2.5), SLIDE_WIDTH - Inches(2), Inches(3))
    p = box.text_frame.add_paragraph()
    p.text = "无论你选择硬件的霸道,还是软件的智慧,\n这都是一个记忆永驻的新时代。"
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER

    # --- Page 30: Closing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
    p = box.text_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    sub = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "The Architecture of Memory"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=False)

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V5_30pages.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_ppt_v5()
