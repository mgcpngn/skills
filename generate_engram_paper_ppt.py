"""
DeepSeek Engram Paper - 30 Pages PPT
基于 leijunskill v2.0 §8 论文解构流程生成
应用三层解构、雷军化表达、配图增强
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_ARCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/engram_architecture_1768541436038.png"
IMG_HASH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ngram_hashing_1768541452540.png"
IMG_USHAPE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ushape_scaling_1768541469755.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/prefetch_strategy_1768541489525.png"
IMG_ZIPF = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/zipf_cache_1768541507122.png"

# --- CANVAS (§2.8.1) ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)

# --- LAYOUT ZONES (§2.8.2) ---
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(1.0)
CONTENT_TOP = Inches(1.8)
CONTENT_HEIGHT = Inches(4.0)
FOOTER_TOP = Inches(6.3)

# --- FONT SIZES (§2.4) ---
TITLE_SIZE = Pt(48)
BIG_NUMBER_SIZE = Pt(80)
EMPHASIS_SIZE = Pt(54)
SUBTITLE_SIZE = Pt(28)
BODY_SIZE = Pt(24)

# Colors
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(255, 215, 0)
LIGHT_GRAY = RGBColor(220, 220, 220)
TIANYI_BLUE = RGBColor(0, 191, 255)
DEEPSEEK_GREEN = RGBColor(46, 204, 113)


def create_engram_paper_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    def add_signature(slide, is_cover=True):
        box = slide.shapes.add_textbox(
            (SLIDE_WIDTH - Inches(5)) / 2, FOOTER_TOP,
            Inches(5), Inches(0.9)
        )
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = TIANYI_BLUE
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(22)
        p2.font.color.rgb = GOLD
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def content_slide(title, bullets, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0.5), TITLE_TOP, Inches(12.333), TITLE_HEIGHT)
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        content_box = s.shapes.add_textbox(Inches(0.8), CONTENT_TOP, Inches(11.733), CONTENT_HEIGHT)
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(12)
            if line.startswith("##"):
                para.text = line[2:].strip()
                para.font.size = EMPHASIS_SIZE
                para.font.color.rgb = GOLD
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = SUBTITLE_SIZE
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = "• " + line
                para.font.size = BODY_SIZE
                para.font.color.rgb = LIGHT_GRAY
        add_notes(s, notes)
        return s

    def data_slide(title, number, unit, explanation, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0.5), TITLE_TOP, Inches(12.333), Inches(0.8))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        num_box = s.shapes.add_textbox(Inches(0), Inches(1.8), SLIDE_WIDTH, Inches(1.8))
        p = num_box.text_frame.add_paragraph()
        p.text = number
        p.font.size = BIG_NUMBER_SIZE
        p.font.bold = True
        p.font.color.rgb = DEEPSEEK_GREEN
        p.alignment = PP_ALIGN.CENTER
        unit_box = s.shapes.add_textbox(Inches(0), Inches(3.8), SLIDE_WIDTH, Inches(0.6))
        p = unit_box.text_frame.add_paragraph()
        p.text = unit
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        exp_box = s.shapes.add_textbox(Inches(1), Inches(4.6), SLIDE_WIDTH - Inches(2), Inches(0.8))
        p = exp_box.text_frame.add_paragraph()
        p.text = explanation
        p.font.size = Pt(22)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_notes(s, notes)
        return s

    def split_slide(title, bullets, img_path, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0.5), TITLE_TOP, Inches(12.333), TITLE_HEIGHT)
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        left_box = s.shapes.add_textbox(Inches(0.5), CONTENT_TOP, Inches(5.5), CONTENT_HEIGHT)
        tf = left_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(10)
            if line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = Pt(26)
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = line
                para.font.size = Pt(20)
                para.font.color.rgb = LIGHT_GRAY
        if img_path and os.path.exists(img_path):
            s.shapes.add_picture(img_path, Inches(6.5), CONTENT_TOP, width=Inches(6.0))
        add_notes(s, notes)
        return s

    def closing_slide(title, subtitle, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        sub = s.shapes.add_textbox(Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.8))
        p = sub.text_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(26)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_signature(s, is_cover=False)
        add_notes(s, notes)
        return s

    # ==================== 30 PAGES ====================

    # --- P1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "DeepSeek Engram"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "条件记忆：稀疏性的新维度"
    p2.font.size = Pt(32)
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.CENTER
    if os.path.exists(IMG_ARCH):
        s.shapes.add_picture(IMG_ARCH, (SLIDE_WIDTH - Inches(5)) / 2, Inches(2.0), width=Inches(5))
    sub = s.shapes.add_textbox(Inches(0), Inches(5.3), SLIDE_WIDTH, Inches(0.5))
    p = sub.text_frame.add_paragraph()
    p.text = "arXiv 2601.07372 | DeepSeek-AI & 北京大学"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)
    add_notes(s, """朋友们，大家好！

今天，我要给大家介绍一篇可能改变AI未来的论文——DeepSeek Engram。

这篇论文问了一个简单但深刻的问题："AI为什么要浪费算力去'回忆'那些早已知道的事实？"

传统大模型就像一个天才学霸，什么都记在脑子里。但Engram说：为什么不带本参考书呢？

这就是条件记忆——稀疏性的新维度。""")

    # --- P2: Agenda ---
    content_slide("分享大纲", [
        "1. 问题：计算与检索的二元性",
        "2. 创新：Engram 条件记忆模块",
        "3. 技术：N-gram Hashing + 门控融合",
        "4. 发现：U形 Scaling Law",
        "5. 效果：全面超越纯 MoE 基线",
        "6. 意义：下一代稀疏模型的范式"
    ], """我们今天的分享分为六个部分。

首先，我们看看当前大模型存在什么问题——为什么它们既浪费算力又记不住东西。

然后，我来揭示Engram的核心创新——条件记忆模块。

接着，我们深入技术细节：N-gram Hashing 和 上下文门控如何工作。

之后，分享一个重要发现——U形 Scaling Law。

再然后，用数据说话，展示Engram的实验效果。

最后，我们一起思考：这对AI的未来意味着什么。""")

    # --- P3: Problem 1 ---
    content_slide("问题：语言模型的二元性", [
        "#组合推理 (Compositional Reasoning)",
        "需要深度、动态的计算",
        "#知识检索 (Knowledge Retrieval)",
        "局部、静态、高度模式化",
        "## 一个需要计算，一个适合查表"
    ], """朋友们，让我们先理解一个根本性的问题。

语言建模其实包含两种完全不同的任务：

第一种是"组合推理"——比如解决数学题、写复杂代码。这需要深度思考，需要强大的计算能力。

第二种是"知识检索"——比如回答"北京是哪国首都"、"1+1等于几"。这些是静态的事实，模式化程度很高。

问题是：现在的大模型，把这两件事混在一起做！用同样的神经网络既做推理，又做检索。

这就好比让爱因斯坦每天花8小时背九九乘法表，太浪费了！""")

    # --- P4: Problem 2 ---
    content_slide("问题：Transformer 的低效", [
        "缺乏原生的知识查表机制",
        "用计算来模拟检索",
        "静态知识在早期层反复重建",
        "## 浪费深度，浪费算力"
    ], """更深入地看，Transformer 架构有一个根本性的缺陷：

它没有"查表"的原生能力！

想象一下，每次你问AI"北京是哪国首都"，它不是直接告诉你，而是：
- 第一层：思考"这是个地理问题"
- 第二层：回想"亚洲有哪些国家"
- 第三层：定位"中国的首都"
- 第四层：输出"北京"

四层计算，就为了回答一个常识问题！

论文的数据显示，模型需要消耗多个早期层来"重建"这些静态的查表操作。这些深度本可以用于更复杂的推理！""")

    # --- P5: Core Innovation ---
    content_slide("核心创新：条件记忆", [
        "#Conditional Computation (MoE)",
        "动态选择用哪些专家计算",
        "#Conditional Memory (Engram)",
        "静态查询预存的知识嵌入",
        "## 计算与记忆，双轴并行"
    ], """DeepSeek的核心创新是什么？

他们提出了一个新概念：条件记忆 (Conditional Memory)。

过去，我们有条件计算——就是MoE，根据输入动态选择用哪些专家来处理。

现在，我们增加了条件记忆——根据输入查询预存的知识嵌入表。

两者是什么关系？

MoE是"选哪个专家来想答案"；
Engram是"选哪本参考书来翻阅"。

这不是替代关系，而是互补关系！就像团队既需要专家动脑子，也需要资料库做参考。""")

    # --- P6: Engram Architecture ---
    split_slide("Engram 架构总览", [
        "#双阶段处理",
        "1. 检索：N-gram → Hash → 嵌入",
        "2. 融合：门控 + 卷积 + 残差",
        "#嵌入模块外挂Transformer",
        "不改变原有注意力和MoE",
        "#仅在特定层应用",
        "平衡性能与延迟"
    ], IMG_ARCH,
    """让我们来看Engram的架构。

它有两个核心阶段：

第一阶段是检索。输入的Token序列先做N-gram提取，然后通过Hash函数转成索引，最后从巨大的嵌入表中查出对应的知识向量。

第二阶段是融合。查到的知识向量不是直接拼接，而是通过一个智能门控机制，和当前的隐藏状态融合，再经过卷积增强，最后通过残差连接回到主干网络。

关键点：Engram是外挂的！它不改变原有的注意力和MoE结构，只是在特定层加入记忆查询。

这种设计让它可以无缝集成到现有模型中。""")

    # --- P7: Tokenizer Compression ---
    data_slide("词表压缩", "23%", "有效词汇量减少", "通过规范化标识符实现",
    """Engram的第一个技术创新：词表压缩。

传统分词器会给"Apple"和"apple"分配不同的ID。但它们其实是同一个概念！

通过NFKC规范化、小写转换等技术，Engram把128k词表压缩了23%。

这意味着什么？

原来需要100亿个词条的知识表，现在只需要77亿。存储省了，查询也快了。

这就是工程师思维——不放过任何一个优化点！""")

    # --- P8: N-gram Hashing ---
    split_slide("N-gram Hashing", [
        "#多头哈希 (Multi-head)",
        "K个独立的哈希函数",
        "减少碰撞，增强鲁棒性",
        "#O(1) 复杂度",
        "无论表多大，查询恒定时间",
        "#确定性索引",
        "支持运行时预取"
    ], IMG_HASH,
    """（三层解构：比喻转化）

N-gram Hashing是Engram的核心技术。

想象你有一本1000亿词条的超级字典。现在我说"北京"，你要找到北京的解释。

如果从头翻到尾？翻到猴年马月也翻不完。

聪明的办法：给每个词条一个编号，按编号存放。

我说"北京"→你算出编号"12345"→直接翻到第12345页。

这就是O(1)复杂度！无论字典多大，查询时间都是恒定的。

而且是多头哈希——用K个不同的哈希函数，减少碰撞，就像用多个索引来交叉验证。""")

    # --- P9: Context-aware Gating ---
    content_slide("上下文感知门控", [
        "查到的嵌入可能有噪声",
        "哈希碰撞/多义词歧义",
        "#解决方案：注意力门控",
        "当前隐藏状态作为Query",
        "检索记忆作为Key/Value",
        "## 语义不匹配→门控关闭"
    ], """查到的知识向量不能直接用，为什么？

首先，哈希可能有碰撞——两个不同的词可能算出同一个索引。

其次，一词多义——"苹果"可能指水果，也可能指公司。

怎么解决？上下文门控！

当前的隐藏状态h已经通过注意力层聚合了全局上下文。把h作为Query，查到的记忆e作为Key和Value，做一个轻量级的注意力。

如果语义匹配，门控打开，知识融入；
如果语义冲突，门控关闭，噪声被抑制。

这个设计非常优雅——用上下文来过滤噪声！""")

    # --- P10: Multi-branch Integration ---
    content_slide("多分支集成", [
        "现代架构使用多分支残差流",
        "#参数共享策略",
        "共享：嵌入表 + Value投影",
        "独立：每个分支的Key投影",
        "## FP8 矩阵乘融合"
    ], """Engram如何集成到最新的多分支架构？

论文使用了一个巧妙的参数共享策略：

嵌入表只有一份——这是最大的参数量，必须共享。
Value投影矩阵也共享——输出格式统一。
但Key投影矩阵每个分支独立——允许分支特定的门控行为。

这样设计的好处？可以把所有线性投影融合成一个FP8矩阵乘法，充分利用GPU的Tensor Core。

这就是论文强调的"infrastructure-aware"——不仅考虑算法，还要考虑硬件实现效率。""")

    # --- P11: Prefetching ---
    split_slide("确定性预取策略", [
        "#MoE: 动态路由",
        "依赖运行时隐藏状态",
        "无法提前知道选哪个专家",
        "#Engram: 确定性索引",
        "仅依赖输入Token序列",
        "可提前计算Hash ID",
        "## 异步预取，隐藏延迟"
    ], IMG_PREFETCH,
    """（三层解构：生活化比喻）

Engram相比MoE有一个巨大的系统优势：确定性预取！

MoE的动态路由依赖运行时的隐藏状态。你必须算完上一层，才能知道这一层选哪个专家。没法提前准备。

但Engram不同！它的索引只依赖输入Token序列——这在推理开始前就知道了！

这就像点外卖。你知道自己中午要吃饭，不会等到饿了再点。提前下单，等你饿的时候，外卖已经送到了。

Engram同样如此：GPU在算第1层的时候，内存已经在准备第3层需要的数据了。完美重叠。

结果？100B参数表卸载到主机内存，吞吐损失不到3%！""")

    # --- P12: Zipf Distribution ---
    split_slide("Zipf 分布与多级缓存", [
        "#自然语言的统计规律",
        "少数N-gram高频出现",
        "大多数N-gram很少出现",
        "#多级存储策略",
        "热门：GPU HBM / Host DRAM",
        "冷门：NVMe SSD"
    ], IMG_ZIPF,
    """自然语言有一个著名的统计规律——Zipf定律。

少数词语出现频率极高，大多数词语很少出现。

就像你手机里的App。微信、抖音每天用，放在首屏。某个三年前下的App从来不用，藏在最后一屏。

Engram利用这个规律设计了多级缓存：

热门的N-gram嵌入→放在GPU显存或主机DRAM，快速访问。
冷门的N-gram嵌入→放在NVMe SSD，容量大但稍慢。

这样，既保证了速度，又实现了超大规模扩展。

这就是工程的智慧——顺应自然规律，而不是逆天而行！""")

    # --- P13: U-shaped Law Discovery ---
    split_slide("U形 Scaling Law", [
        "#稀疏分配问题",
        "固定参数预算，如何分配?",
        "#实验发现",
        "纯MoE (ρ=100%) 不是最优",
        "纯Engram (ρ=0%) 也不是最优",
        "## 最优点在 ρ ≈ 75-80%"
    ], IMG_USHAPE,
    """这是论文最重要的发现之一——U形 Scaling Law！

研究团队问了一个关键问题：如果总参数预算固定，应该给MoE分多少，给Engram分多少？

他们做了大量实验，结果令人惊讶：

纯MoE（ρ=100%）：模型缺乏专用记忆，被迫用深度来重建静态知识。
纯Engram（ρ=0%）：模型失去动态计算能力，无法处理需要推理的任务。

最优解在中间！大约75-80%给MoE，20-25%给Engram。

这证明了两者是结构性互补——计算不能替代记忆，记忆也不能替代计算！""")

    # --- P14: Key Metrics ---
    data_slide("最优分配比例", "75-80%", "给 MoE 专家", "其余 20-25% 给 Engram 记忆",
    """让我强调一下这个数字：75-80%！

这意味着：在固定预算下，你应该把大约四分之三的稀疏参数给MoE专家，四分之一给Engram记忆。

而且这个比例在不同规模下都很稳定——无论是5B模型还是10B模型，最优点都在这个区间。

这是一个可操作的工程指导原则。如果你在设计下一代稀疏模型，记住这个数字！""")

    # --- P15: Experimental Setup ---
    content_slide("实验配置", [
        "#训练数据: 262B Tokens",
        "#分词器: DeepSeek-v3 (128k)",
        "#架构: 30层Transformer + MLA",
        "#激活参数: 固定 3.8B",
        "## 严格对照实验"
    ], """让我介绍一下实验配置。

所有模型使用完全相同的训练数据——262B tokens。
使用DeepSeek-v3的128k词表分词器。
架构是30层Transformer，使用多头潜在注意力（MLA）。

关键是：所有模型的激活参数都固定在3.8B！

这是真正的公平对比——相同的计算量，相同的训练数据，只有架构不同。

任何性能差异，都来自架构设计本身。""")

    # --- P16: Model Comparison ---
    content_slide("模型对比", [
        "#Dense-4B: 4.1B 参数",
        "密集模型基线",
        "#MoE-27B: 26.7B 参数",
        "72个路由专家 + 2个共享专家",
        "#Engram-27B: 26.7B 参数",
        "55个路由专家 + 5.7B Engram"
    ], """论文训练了四个模型：

Dense-4B：密集模型基线，4.1B参数，作为对照。

MoE-27B：纯MoE架构，26.7B参数，72个路由专家。

Engram-27B：从MoE-27B派生，把17个专家的参数（5.7B）重新分配给Engram记忆模块。总参数量不变，还是26.7B。

Engram-40B：进一步扩大Engram记忆到18.5B，总参数达到39.5B。

关键是Engram-27B和MoE-27B的对比——严格等参数、等FLOPs！""")

    # --- P17: MMLU Performance ---
    data_slide("知识问答：MMLU", "+3.4", "相比 MoE-27B", "57.4 → 60.4 准确率",
    """让数据说话！

在MMLU知识问答测试中，Engram-27B比MoE-27B高出3.4分！

从57.4%提升到60.4%。

这完全符合预期——MMLU测试的是知识记忆能力，正是Engram的主场。

但别急，更惊人的结果还在后面……""")

    # --- P18: BBH Performance ---
    data_slide("复杂推理：BBH", "+5.0", "相比 MoE-27B", "50.9 → 55.9 准确率",
    """在BBH（Big Bench Hard）复杂推理测试中，Engram-27B比MoE-27B高出5分！

等等，这是推理任务啊！Engram是记忆模块，为什么能提升推理能力？

论文给出了解释：Engram把静态知识的重建工作卸载了，让早期层从"回忆事实"的负担中解放出来。

这些深度就可以用于更复杂的推理了！

这就像学霸不用再背九九乘法表，可以把精力用在证明定理上。""")

    # --- P19: ARC Performance ---
    data_slide("科学推理：ARC-Challenge", "+3.7", "相比 MoE-27B", "70.1 → 73.8 准确率",
    """ARC-Challenge科学推理测试：+3.7分！

从70.1%提升到73.8%。

科学推理需要什么？既需要记住科学事实，又需要逻辑推导。

Engram两边都帮上忙了——记忆模块负责事实，释放出来的深度负责推理。""")

    # --- P20: Code Performance ---
    data_slide("代码生成：HumanEval", "+3.0", "相比 MoE-27B", "37.8 → 40.8 Pass@1",
    """HumanEval代码生成测试：+3.0分！

为什么代码任务也能提升？

编程有很多固定模式：遍历数组、读写文件、网络请求……这些都是可以查表的！

Engram把常见的代码模式存在记忆表里，需要的时候直接查出来用。

就像有了一本代码片段手册。""")

    # --- P21: Math Performance ---
    data_slide("数学推理：MATH", "+2.4", "相比 MoE-27B", "28.3 → 30.7 准确率",
    """MATH数学推理测试：+2.4分！

数学需要什么？公式记忆 + 逻辑推导。

公式记住了，推导自然更流畅。""")

    # --- P22: Long Context ---
    data_slide("长上下文：Multi-Query NIAH", "97.0", "准确率 (vs 84.2)", "提升 12.8 个百分点",
    """长上下文性能提升更加惊人！

Multi-Query NIAH测试：从84.2%提升到97.0%，提升12.8个百分点！

为什么长上下文提升这么大？

因为Engram把局部依赖关系卸载了，注意力可以专注于全局上下文。

就像不用记琐事，脑子可以专注于大局。""")

    # --- P23: Mechanistic Analysis ---
    content_slide("机制分析", [
        "#LogitLens 分析",
        "早期层不再重建静态知识",
        "有效深度增加",
        "#CKA 相似性分析",
        "层间表示更加独立",
        "## 释放深度用于推理"
    ], """论文还做了深入的机制分析。

通过LogitLens技术观察模型内部，发现：

在MoE模型中，早期层在"重建"静态知识——比如解析"北京是中国首都"。

在Engram模型中，这个工作被记忆模块接管了，早期层可以直接开始更高级的处理。

CKA相似性分析也证实：Engram模型的层间表示更加独立，说明每一层都在做不同的事情。

简单说：Engram增加了模型的"有效深度"！""")

    # --- P24: Training Cost ---
    data_slide("训练成本对比", "0%", "额外算力开销", "等参数、等FLOPs 严格对照",
    """Engram需要额外的训练成本吗？

不需要！

Engram-27B和MoE-27B的激活参数数量完全相同（3.8B），所以每个Token的FLOPs也相同。

嵌入表虽然参数量大，但每个Token只检索常数个slot，不增加计算量。

等参数、等FLOPs，纯粹是免费的性能提升！""")

    # --- P25: Inference Efficiency ---
    data_slide("推理效率", "<3%", "100B 卸载的吞吐损失", "确定性预取 + 异步传输",
    """推理时把100B参数的Engram表卸载到主机内存，吞吐损失有多少？

不到3%！

怎么做到的？确定性预取。

因为Hash索引只依赖输入Token，可以在推理开始前就计算好。然后异步从主机内存取，和GPU计算重叠。

便宜的DDR内存，发挥了97%的GPU显存效果。这就是软件工程的力量！""")

    # --- P26: Strategic Value 1 ---
    content_slide("战略价值：AI民主化", [
        "无需顶级GPU",
        "RTX 4090 + 便宜DDR即可",
        "打破显存瓶颈",
        "## 让每个人都用得起大模型"
    ], """Engram的战略价值是什么？

AI民主化！

过去，跑大模型需要昂贵的H100，几十万一块。

有了Engram，你可以把大部分参数放在便宜的DDR内存里。一块RTX 4090加几百块钱内存条，就能跑起来。

打破显存瓶颈，让每个人都用得起大模型。

这就像PC普及让每个人都有电脑一样——Engram让大模型走向每个人。""")

    # --- P27: Strategic Value 2 ---
    content_slide("战略价值：自主可控", [
        "面对芯片出口管制",
        "硬件不够，软件来凑",
        "用算法弥补显存差距",
        "## More with Less"
    ], """还有一层更深的意义——自主可控。

在当前形势下，很多地方拿不到最新GPU。传统方案：增加更多专家→需要更多显存→买不到GPU→完蛋。

Engram方案：增加记忆模块→卸载到便宜DDR→用软件隐藏延迟→性能反超。

用算法智慧弥补硬件差距。这就是"More with Less"的精神！""")

    # --- P28: Future Vision ---
    content_slide("未来展望", [
        "#条件记忆成为必备原语",
        "下一代稀疏模型的标配",
        "#与MoE深度融合",
        "计算+记忆的联合优化",
        "## 稀疏性的新维度"
    ], """论文的结论是什么？

条件记忆将成为下一代稀疏模型的"必备原语"！

就像MoE从可选变成标配一样，Engram也会成为标配。

未来的模型会同时拥有：
- 条件计算（MoE）——动态选择专家
- 条件记忆（Engram）——静态查询知识

两个维度，联合优化。

这是稀疏性的新维度，这是AI架构的新范式！""")

    # --- P29: Summary ---
    content_slide("总结", [
        "#技术创新",
        "条件记忆 + N-gram Hashing + 门控融合",
        "#科学发现",
        "U形 Scaling Law，最优分配 75-80%",
        "#实验验证",
        "BBH +5.0, MMLU +3.4, 长上下文 +12.8"
    ], """让我来做个总结。

技术创新：Engram提出了条件记忆模块，通过N-gram Hashing实现O(1)查表，通过门控机制融合上下文。

科学发现：U形Scaling Law揭示了计算与记忆的结构性互补，最优分配比例75-80%。

实验验证：在多个基准测试中全面领先——BBH +5.0，MMLU +3.4，长上下文 +12.8。

这不是小修小补，这是范式的转变！""")

    # --- P30: Closing ---
    closing_slide("谢谢", "Conditional Memory via Scalable Lookup",
    """朋友们，这就是今天的全部内容。

DeepSeek Engram告诉我们一个简单的道理：不要让AI浪费算力在已知的事情上。

计算归计算，记忆归记忆。各司其职，两全其美。

这是一个优雅的解决方案，也是一个深刻的洞察。

论文的代码已经开源。我相信，很快我们会看到更多基于Engram的创新。

感谢大家的聆听！朋友们，我们下次再见！""")

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'DeepSeek_Engram_Paper_30pages.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_engram_paper_ppt()
