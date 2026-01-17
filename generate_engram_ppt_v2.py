"""
DeepSeek Engram - 20 Pages PPT (V2 - Fixed Layout)
修复布局问题：
1. 封面标题与图片分离
2. 强调文字字号限制在72pt
3. 图文页严格对齐
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_NGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_ngram_hashing_1768530510183.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_prefetch_strategy_1768530526058.png"

# --- CANVAS (leijunskill §2.8.1) ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)

# --- LAYOUT ZONES (leijunskill §2.8.2) ---
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(1.0)
CONTENT_TOP = Inches(1.8)
CONTENT_HEIGHT = Inches(4.0)
FOOTER_TOP = Inches(6.3)

# --- FONT SIZES (leijunskill §2.4 - with max limits) ---
TITLE_SIZE = Pt(48)          # Max: 60pt
BIG_NUMBER_SIZE = Pt(80)     # Max: 96pt
EMPHASIS_SIZE = Pt(54)       # Max: 72pt (for ## lines)
SUBTITLE_SIZE = Pt(28)       # Max: 36pt
BODY_SIZE = Pt(24)           # Max: 28pt

# Colors
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(255, 215, 0)
LIGHT_GRAY = RGBColor(220, 220, 220)
TIANYI_BLUE = RGBColor(0, 191, 255)
DEEPSEEK_GREEN = RGBColor(46, 204, 113)


def create_engram_ppt_v2():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    def add_signature(slide, is_cover=True):
        """Signature at footer zone (§2.8.3)"""
        box = slide.shapes.add_textbox(
            (SLIDE_WIDTH - Inches(5)) / 2,
            FOOTER_TOP,
            Inches(5), Inches(0.9)
        )
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = TIANYI_BLUE
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(22)
        p2.font.color.rgb = GOLD
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def content_slide(title, bullets, notes):
        """Standard content slide (§2.8.4)"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title zone
        box = s.shapes.add_textbox(
            Inches(0.5), TITLE_TOP,
            Inches(12.333), TITLE_HEIGHT
        )
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Content zone
        content_box = s.shapes.add_textbox(
            Inches(0.8), CONTENT_TOP,
            Inches(11.733), CONTENT_HEIGHT
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(12)
            if line.startswith("##"):
                para.text = line[2:].strip()
                para.font.size = EMPHASIS_SIZE  # 54pt, not 120pt!
                para.font.color.rgb = GOLD
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = SUBTITLE_SIZE
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = "• " + line
                para.font.size = BODY_SIZE
                para.font.color.rgb = LIGHT_GRAY
        add_notes(s, notes)
        return s

    def data_slide(title, number, unit, explanation, notes):
        """Big number slide (§2.8.5)"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(
            Inches(0.5), TITLE_TOP,
            Inches(12.333), Inches(0.8)
        )
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Big Number
        num_box = s.shapes.add_textbox(
            Inches(0), Inches(1.8),
            SLIDE_WIDTH, Inches(1.8)
        )
        p = num_box.text_frame.add_paragraph()
        p.text = number
        p.font.size = BIG_NUMBER_SIZE  # 80pt
        p.font.bold = True
        p.font.color.rgb = DEEPSEEK_GREEN
        p.alignment = PP_ALIGN.CENTER
        # Unit
        unit_box = s.shapes.add_textbox(
            Inches(0), Inches(3.8),
            SLIDE_WIDTH, Inches(0.6)
        )
        p = unit_box.text_frame.add_paragraph()
        p.text = unit
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Explanation
        exp_box = s.shapes.add_textbox(
            Inches(1), Inches(4.6),
            SLIDE_WIDTH - Inches(2), Inches(0.8)
        )
        p = exp_box.text_frame.add_paragraph()
        p.text = explanation
        p.font.size = Pt(22)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_notes(s, notes)
        return s

    def split_slide(title, bullets, img_path, notes):
        """Left text, right image (§2.8.6) - FIXED alignment"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(
            Inches(0.5), TITLE_TOP,
            Inches(12.333), TITLE_HEIGHT
        )
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Left content (start at CONTENT_TOP)
        left_box = s.shapes.add_textbox(
            Inches(0.5), CONTENT_TOP,
            Inches(5.5), CONTENT_HEIGHT
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(10)
            if line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = Pt(26)
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = line
                para.font.size = Pt(20)
                para.font.color.rgb = LIGHT_GRAY
        # Right image (SAME top as text for alignment!)
        if img_path and os.path.exists(img_path):
            s.shapes.add_picture(
                img_path,
                Inches(6.5),        # Left edge at 6.5"
                CONTENT_TOP,        # SAME top as text!
                width=Inches(6.0),  # Width 6"
            )
        add_notes(s, notes)
        return s

    def closing_slide(title, subtitle, notes):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        sub = s.shapes.add_textbox(Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.8))
        p = sub.text_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(26)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_signature(s, is_cover=False)
        add_notes(s, notes)
        return s

    # ==================== 20 PAGES: ENGRAM DEEP DIVE ====================

    # --- P1: Cover (FIXED: title and image separated) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # Title in title zone (0.3" - 1.3")
    box = s.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1.0)
    )
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "DeepSeek Engram"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "条件记忆：稀疏性的新维度"
    p2.font.size = Pt(32)
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.CENTER
    # Image in content zone (2.0" start)
    if os.path.exists(IMG_ENGRAM):
        img_width = Inches(5)
        img_left = (SLIDE_WIDTH - img_width) / 2
        s.shapes.add_picture(IMG_ENGRAM, img_left, Inches(2.0), width=img_width)
    # Source text below image
    sub = s.shapes.add_textbox(Inches(0), Inches(5.3), SLIDE_WIDTH, Inches(0.5))
    p = sub.text_frame.add_paragraph()
    p.text = "Conditional Memory via Scalable Lookup | arXiv 2601.07372"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)
    add_notes(s, """朋友们，大家好！

今天，我想给大家介绍一项可能改变AI未来的技术——DeepSeek Engram。

这不是普通的技术改进，而是一个全新范式。Engram问了一个简单却深刻的问题："AI为什么要浪费算力去'回忆'那些早已知道的事实？"

这就是今天的主角——条件记忆，稀疏性的新维度。""")

    # --- P2: Agenda ---
    content_slide("分享大纲", [
        "1. 问题：AI的失忆与浪费",
        "2. 洞察：重构 vs. 查表",
        "3. 方案：Engram 技术架构",
        "4. 效果：性能与效率双赢",
        "5. 意义：AI的未来范式"
    ], """我们今天的分享分为五个部分。首先看问题，然后揭示洞察，接着深入技术，之后用数据说话，最后思考意义。准备好了吗？Let's go！""")

    # --- P3: AI Amnesia ---
    content_slide("问题：AI的失忆症", [
        "上下文窗口有限",
        "KV Cache 指数膨胀",
        "聊着聊着就忘了",
        "## AI 记不住对话"
    ], """朋友们，我相信你们都遇到过这种情况：和AI聊了半个小时，突然问它"我刚才说的那个想法怎样？"，它愣住了说"抱歉，你说过什么想法吗？"这就是AI的失忆症。""")

    # --- P4: Computational Waste ---
    content_slide("问题：算力的浪费", [
        "90% 计算在重构静态知识",
        "北京是首都——每次都要重新推导",
        "1+1=2——每次都要重新计算",
        "## 这是巨大的浪费"
    ], """更糟糕的是，AI还在浪费大量算力。研究表明，90%的计算都花在"重构"那些早已知道的静态知识上。"北京是首都"每次都要重新推导，不荒谬吗？""")

    # --- P5: Core Insight ---
    content_slide("核心洞察", [
        "#重构 (Reconstruction)",
        "用神经网络计算出答案",
        "#查表 (Lookup)",
        "直接从记忆中取出答案",
        "## 别算了，直接查！"
    ], """DeepSeek的洞察很简单：为什么不直接查表呢？就像考试，天才学霸什么都记在脑子里，但聪明学生带参考书，脑子记不住的翻书就行。这就是Engram的精髓：别算了，直接查！""")

    # --- P6: What is Engram (FIXED: image aligned) ---
    split_slide("什么是 Engram", [
        "#条件记忆模块",
        "Conditional Memory Module",
        "#稀疏性的新维度",
        "与 MoE 互补，而非替代",
        "#O(1) 复杂度查表",
        "无论表多大，查询恒定时间"
    ], IMG_ENGRAM,
    """Engram是一个条件记忆模块，嵌入到Transformer架构中。它代表了新的稀疏性维度——MoE解决"用哪些神经元计算"，Engram解决"查哪些记忆"。关键是O(1)复杂度！""")

    # --- P7: N-gram Hashing ---
    split_slide("技术：N-gram Hashing", [
        "#输入 Token 编码为 Hash ID",
        "确定性、上下文相关",
        "#用 Hash 作为索引",
        "直接定位知识向量",
        "#多头哈希",
        "捕捉不同层次的模式"
    ], IMG_NGRAM,
    """核心技术叫N-gram Hashing。把输入Token序列转成确定性的数字，用这个数字去查表。就像给字典编号，说"北京"就算出编号12345，直接翻到那页。快得让人难以置信。""")

    # --- P8: Tokenizer Compression ---
    data_slide("词表压缩", "23%", "有效词汇量减少", "通过规范化标识符实现",
    """Engram还把有效词汇量减少了23%！"AI"、"A.I."、"artificial intelligence"其实指同一个东西，通过规范化，省了23%存储，查询效率也更高。""")

    # --- P9: Memory Hierarchy ---
    content_slide("多级缓存架构", [
        "#热门词条 → GPU HBM",
        "高频访问，驻留显存",
        "#长尾词条 → 主机 DDR",
        "低频访问，按需加载",
        "## 顺应 Zipf 定律"
    ], """Engram利用Zipf定律：少数词语出现频率极高。热门词条放GPU显存，长尾词条放DDR内存。既保证速度，又节省昂贵显存。顺应自然规律，这就是工程的智慧。""")

    # --- P10: Prefetching ---
    split_slide("确定性预取", [
        "#计算先于需求",
        "提前算出 Hash ID",
        "#异步加载数据",
        "GPU计算时，内存已在准备",
        "#吞吐损失 < 3%",
        "几乎无感知"
    ], IMG_PREFETCH,
    """从DDR取数据比GPU内部慢，怎么办？提前预取！因为Hash是确定性的，可以提前算出需要哪些数据。GPU算第1层时，内存已在准备第3层的数据。结果：吞吐损失不到3%！""")

    # --- P11: Key Numbers ---
    data_slide("核心参数", "100B", "卸载到主机内存", "27B 活跃模型 + 100B 知识表",
    """100B参数的知识表，全部存在DDR内存中！活跃模型只有27B。27B+100B=127B，但只有27B在GPU上。传统方案需要几百GB显存，Engram只要27B的GPU需求加几百块DDR。""")

    # --- P12: BBH Performance ---
    data_slide("性能：复杂推理", "+5.0", "BBH 测试提升", "vs 同等 FLOPs 的 MoE 模型",
    """在BBH复杂推理测试中，Engram比同规模MoE模型高出5分！同样的算力预算，Engram全面碾压。因为它把静态知识回忆的工作卸载了，神经网络可以专注于真正的推理。""")

    # --- P13: MMLU Performance ---
    data_slide("性能：知识问答", "+3.4", "MMLU 测试提升", "知识检索准确率显著提高",
    """在MMLU知识问答测试中提升3.4分。这完全符合预期——知识问答正是Engram的主场。参数多不等于记得准，Engram的外挂知识表存储精度比神经网络自己"悟"出来的更高。""")

    # --- P14: Code Performance ---
    data_slide("性能：代码生成", "+3.0", "HumanEval 测试提升", "常见编程模式召回更准确",
    """HumanEval代码生成测试提升3.0分。编程有很多固定模式：遍历数组、读写文件、网络请求……Engram把这些模式存在知识表里，需要时查出来用。让AI也有了老手的经验库。""")

    # --- P15: U-shaped Law ---
    content_slide("U形 Scaling Law", [
        "#稀疏分配问题",
        "MoE vs Engram 如何分配?",
        "#发现 U形曲线",
        "纯MoE和纯Engram都不是最优",
        "## 混合分配最佳"
    ], """DeepSeek发现了U形Scaling Law：纯MoE或纯Engram都不是最优。最优是混合分配——约20-25%给Engram，其余给MoE。就像团队既要思考者也要执行者。这是用数据说话的科学。""")

    # --- P16: AI Democratization ---
    content_slide("战略价值：AI民主化", [
        "无需顶级 GPU",
        "RTX 4090 + DDR 即可",
        "打破 NVIDIA 垄断",
        "## 让每个人都用得起 AI"
    ], """Engram代表AI的民主化。过去跑大模型需要价值几十万的H100，有了Engram，一块RTX 4090加几百块DDR就能跑GPT-4级效果。这像当年PC的普及——让AI从云端走向每个人的桌面。""")

    # --- P17: Self-reliance ---
    content_slide("战略价值：自主可控", [
        "面对芯片出口管制",
        "软件弥补硬件差距",
        "DeepSeek V3 仅用 $5.58M 训练",
        "## More with Less"
    ], """在出口管制背景下，很多地方拿不到最新GPU。Engram告诉我们：硬件不够，软件来凑。DeepSeek V3训练成本558万美元，GPT-4据说1亿。同样结果，成本只有1/20。以小博大！""")

    # --- P18: Agentic AI ---
    content_slide("未来：Agentic AI", [
        "从 Chatbot 到 Agent",
        "长期记忆成为刚需",
        "Engram = AI的长期记忆引擎",
        "## 记忆，是Agent的灵魂"
    ], """AI正在从聊天机器人进化成智能代理。Agent需要记住你的偏好、之前说过的话、任务上下文。没有记忆就没有Agent。Engram正是解决这个问题的关键——AI的长期记忆引擎。""")

    # --- P19: Summary ---
    content_slide("总结：Engram 的价值", [
        "#技术突破",
        "从重构到查表，O(1) 复杂度",
        "#性能提升",
        "BBH +5.0，MMLU +3.4",
        "#战略意义",
        "AI民主化，自主可控"
    ], """总结Engram的价值：技术上实现从重构到查表的范式转变；性能上多项测试全面领先；战略上代表AI民主化和自主可控。这不是某个天才的成果，而是整个团队日夜奋战的结晶。""")

    # --- P20: Closing ---
    closing_slide("谢谢", "Conditional Memory via Scalable Lookup",
    """朋友们，这就是今天的全部内容。Engram不仅是技术模块，更代表新的思维方式：不要让AI浪费算力在已知的事情上。这只是开始，未来会有更多基于Engram的创新。感谢聆听！""")

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'DeepSeek_Engram_20pages_v2.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_engram_ppt_v2()
