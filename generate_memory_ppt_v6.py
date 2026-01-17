"""
Memory Architecture Report - V6 (40 Pages)
Generated following leijunskill standards:
- 40 pages as requested
- Technical concept visualizations (§8)
- Signature only on Cover and Closing pages (§5)
- Strict layout rules (§6)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"
# New concept explainer images (§8)
IMG_MEMORY_WALL = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_memory_wall_1768530473358.png"
IMG_KV_CACHE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_kv_cache_1768530494022.png"
IMG_NGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_ngram_hashing_1768530510183.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_prefetch_strategy_1768530526058.png"

# --- CONSTANTS (leijunskill §6.1) ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)

# --- SAFE ZONES (leijunskill §6.2) ---
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(1.5)
CONTENT_TOP = Inches(2.0)
CONTENT_HEIGHT = Inches(4.0)

# --- SPLIT LAYOUT (leijunskill §6.3) ---
SPLIT_TEXT_LEFT = Inches(0.5)
SPLIT_TEXT_WIDTH = Inches(5.5)
SPLIT_IMG_LEFT = Inches(6.5)
SPLIT_IMG_WIDTH = Inches(6.3)


def create_ppt_v6():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_signature(slide, is_cover=True):
        """Stylized signature - leijunskill §5 - ONLY for cover and closing"""
        sig_width = Inches(5)
        sig_height = Inches(0.9)
        sig_left = (SLIDE_WIDTH - sig_width) / 2
        sig_top = SLIDE_HEIGHT - sig_height - Inches(0.15)
        box = slide.shapes.add_textbox(sig_left, sig_top, sig_width, sig_height)
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(0, 191, 255)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(255, 215, 0)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def add_title(slide, text, center=False):
        box = slide.shapes.add_textbox(Inches(0.5), TITLE_TOP, SLIDE_WIDTH - Inches(1), TITLE_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if center:
            p.alignment = PP_ALIGN.CENTER

    def add_content(slide, lines):
        box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, SLIDE_WIDTH - Inches(2), CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(16)
            if line.startswith("##"):
                p.text = line[2:].strip()
                p.font.size = Pt(54)
                p.font.color.rgb = RGBColor(255, 215, 0)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = "• " + line
                p.font.size = Pt(26)
                p.font.color.rgb = RGBColor(220, 220, 220)

    def add_split(slide, lines, img_path):
        box = slide.shapes.add_textbox(SPLIT_TEXT_LEFT, CONTENT_TOP, SPLIT_TEXT_WIDTH, CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(12)
            if line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(220, 220, 220)
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, SPLIT_IMG_LEFT, CONTENT_TOP, width=SPLIT_IMG_WIDTH)

    # ========== PAGE GENERATION (40 pages) ==========

    # --- P1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "内存的战争\nThe Architecture of Memory", center=True)
    if os.path.exists(IMG_TITLE):
        s.shapes.add_picture(IMG_TITLE, Inches(4), Inches(2.2), height=Inches(3.0))
    sub = s.shapes.add_textbox(Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "NVIDIA Rubin vs. DeepSeek Engram | 2026 内存架构深度解析"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)

    # --- P2: Agenda ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "汇报大纲")
    add_content(s, [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件霸权",
        "3. DeepSeek Engram：软件革命",
        "4. 路线对比：硬件 vs 软件",
        "5. 性能与经济分析",
        "6. 未来展望：Agentic AI"
    ])

    # --- P3-6: Background (4 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2026：范式转移")
    add_content(s, [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Memory Wall：内存墙")
    add_split(s, [
        "#GPU 算力远超内存带宽",
        "数据传输成为瓶颈",
        "权重访问速度限制推理",
        "#算力越强，内存墙越高"
    ], IMG_MEMORY_WALL)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Context Wall：上下文墙")
    add_content(s, [
        "KV Cache 随序列长度指数膨胀",
        "模型无法记住完整对话",
        "处理百万 Token 需天价计算",
        "## 失忆，是 AI 最大的痛点"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "两大阵营的回答")
    add_content(s, [
        "#NVIDIA (硬件派)",
        "Rubin 架构：极致协同，专有互联",
        "#DeepSeek (软件派)",
        "Engram 技术：算法优化，软件定义内存"
    ])

    # --- P7-16: NVIDIA Rubin (10 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVIDIA Rubin 架构")
    add_split(s, [
        "#Jensen Huang @ CES 2026",
        "年度节奏发布",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "六大芯片生态系统")
    add_content(s, [
        "Rubin GPU - 核心计算引擎",
        "Vera CPU - 内存扩展",
        "NVLink 6 Switch - 高速互联",
        "ConnectX-9 SuperNIC - 网络加速",
        "BlueField-4 DPU - 上下文管理",
        "Spectrum-6 Ethernet Switch"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin GPU：性能怪兽")
    add_content(s, [
        "## 3360 亿",
        "晶体管数量 (台积电 3nm)",
        "## 50 PFLOPS",
        "FP4 推理性能"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：速度的革命")
    add_content(s, [
        "## 22 TB/s",
        "内存带宽 (vs Blackwell 8 TB/s)",
        "## 2.75x",
        "代际提升"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：容量震撼")
    add_content(s, [
        "## 288 GB",
        "单卡 HBM4 显存容量",
        "可容纳 1440 亿参数模型权重",
        "MoE 路由决策可在微秒内完成"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Vera CPU：AI 工厂专属")
    add_content(s, [
        "88 颗 Olympus 定制 ARM 核心",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "首款支持 FP8 的 CPU"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVLink-C2C：高速一致性")
    add_content(s, [
        "## 1.8 TB/s",
        "芯片间互联带宽",
        "相比 Grace-Blackwell 提升 2x",
        "GPU 可像访问本地内存一样访问 CPU 内存"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVL72：超级机柜")
    add_content(s, [
        "72 块 Rubin GPU + 36 块 Vera CPU",
        "## 54 TB",
        "机柜总内存容量",
        "统一内存寻址，透明访问"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "BlueField-4 与 ICMS")
    add_split(s, [
        "#G3.5 上下文存储层",
        "介于 HBM/DRAM 与 SSD 之间",
        "每 GPU 最高 16 TB 上下文",
        "KV Cache 成为共享资源"
    ], IMG_KV_CACHE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin 硬件总结")
    add_content(s, [
        "#核心理念：物理统一内存空间",
        "22 TB/s HBM4 + 1.8 TB/s C2C + 16 TB ICMS",
        "#代价：全栈锁定 NVIDIA 生态",
        "高 CapEx，高性能"
    ])

    # --- P17-24: DeepSeek Engram (8 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "DeepSeek Engram")
    add_split(s, [
        "#软件定义的外挂内存",
        "无需专有硬件",
        "利用主机 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "核心洞察")
    add_content(s, [
        "LLM 大量计算浪费在重构静态知识",
        "为什么不直接查找?",
        "## Reconstruct vs. Lookup",
        "用查表替代计算"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "N-gram Hashing 原理")
    add_split(s, [
        "#输入 Token 序列编码为 Hash",
        "确定性、上下文相关的 ID",
        "用 Hash 作为 Key 查询",
        "## O(1) 复杂度查表"
    ], IMG_NGRAM)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Embedding Table：知识仓库")
    add_content(s, [
        "## 100B 参数",
        "存放在主机 DDR 内存中",
        "27B 活跃模型 + 100B 知识表",
        "小模型拥有万亿知识"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "克服 PCIe 瓶颈")
    add_split(s, [
        "#PCIe 带宽远低于 NVLink",
        "解决方案：确定性预取",
        "提前计算 Hash ID",
        "异步获取数据"
    ], IMG_PREFETCH)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "预取策略效果")
    add_content(s, [
        "在 GPU 执行早期层时",
        "异步从 DDR 获取后续所需数据",
        "## 吞吐量损失 < 3%",
        "廉价 DDR 成为 HBM 的可行扩展"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram vs RAG")
    add_content(s, [
        "#传统 RAG：文档级检索",
        "在 Transformer 外部操作",
        "#Engram：Token 级查表",
        "在 Transformer 内部融合",
        "粒度更细，延迟更低"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram 软件总结")
    add_content(s, [
        "#核心理念：算法替代硬件",
        "N-gram Hashing + 预取重叠",
        "#代价：需修改模型架构",
        "低成本，高可及性"
    ])

    # --- P25-30: Comparison (6 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "路线对比")
    add_split(s, [
        "#NVIDIA: Cohesion (内聚)",
        "硬件透明，模型无感知",
        "#DeepSeek: Retrieval (检索)",
        "软件显式，模型需适配"
    ], IMG_COMPARE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "目标记忆类型对比")
    add_content(s, [
        "#Rubin: 动态 KV Cache",
        "对话上下文 (短期记忆)",
        "每次 Token 生成都需重读",
        "#Engram: 静态 Knowledge",
        "世界知识 (长期记忆)",
        "训练时固化，推理时查表"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "技术参数对比")
    add_content(s, [
        "#NVIDIA Rubin",
        "22 TB/s 带宽, 54 TB 容量, 专有硬件",
        "#DeepSeek Engram",
        "PCIe 带宽, 无限 DDR, 开源软件"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "适用场景对比")
    add_content(s, [
        "#Rubin 适合",
        "超大规模云厂商 (Microsoft/AWS)",
        "需要极致性能和 SLA 保障",
        "#Engram 适合",
        "资源受限环境",
        "消费级硬件部署"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "互补而非替代")
    add_content(s, [
        "Rubin 解决 KV Cache 容量",
        "Engram 解决知识存储",
        "## 两者可以结合",
        "DDR 存知识，HBM 存上下文"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "开源 vs 闭源")
    add_content(s, [
        "#NVIDIA",
        "专有生态，高锁定，高壁垒",
        "#DeepSeek",
        "开源权重，低门槛，普惠创新"
    ])

    # --- P31-34: Performance & Economics (4 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram 性能实测")
    add_content(s, [
        "## +5.0",
        "BBH (复杂推理)",
        "## +3.4",
        "MMLU (知识问答)"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin 经济效益")
    add_content(s, [
        "## 10x",
        "Token 成本下降 (vs Blackwell)",
        "## $5B/$100M",
        "投资回报率 (Token 收入/基础设施投入)"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "AI 民主化")
    add_content(s, [
        "无需 H100/Rubin",
        "RTX 4090 集群跑 GPT-5 级模型",
        "## 打破 NVIDIA 税"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "地缘政治意义")
    add_content(s, [
        "面对出口管制",
        "Engram 是突围方案",
        "DeepSeek V3 仅用 $5.58M 训练",
        "## More with Less"
    ])

    # --- P35-38: Future (4 pages) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "未来：Agentic AI")
    add_content(s, [
        "从 Chatbot 到 Agent",
        "多步推理，长期记忆",
        "管理复杂工作流",
        "## Context is the new Currency"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "百万 Token 上下文")
    add_content(s, [
        "Rubin CPX 可摄入完整代码库",
        "一小时视频作为输入",
        "## 将整个仓库作为工作记忆"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "共享记忆池")
    add_content(s, [
        "BlueField-4 管理 KV Cache",
        "多 Agent 共享，无需重算",
        "## 协作式 AI 成为可能"
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "长期一致性记忆")
    add_content(s, [
        "Engram 支持三种记忆类型",
        "Episodic / Semantic / Procedural",
        "## +15 分 (仅用 1% Token)"
    ])

    # --- P39: Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "总结：内存革命")
    add_content(s, [
        "Rubin 定义了性能上限",
        "Engram 拉高了可及下限",
        "## 殊途同归：解决 AI 失忆"
    ])

    # --- P40: Closing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
    p = box.text_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "The Architecture of Memory"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=False)

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V6_40pages.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_ppt_v6()
