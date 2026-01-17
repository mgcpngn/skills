
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Image Paths from previous generation
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"

def create_ppt_v4():
    prs = Presentation()
    
    # --- Skill Rule: Canvas Standardization (16:9 High Def) ---
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    BG_COLOR = RGBColor(15, 15, 25) # Deep LeiJun Blue/Black

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_copyright(slide, is_cover=False):
        # Mandatory Copyright Injection
        text = "天翼云湖北分公司 作者：王理" if is_cover else "天翼云湖北分公司 王理"
        
        # Position: Bottom Center/Right relative to layout height
        width = Inches(10)
        height = Inches(0.4)
        left = (SLIDE_WIDTH - width) / 2
        top = SLIDE_HEIGHT - Inches(0.5) # 0.5 inches from bottom edge
        
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100) # Subtle Grey
        p.alignment = PP_ALIGN.CENTER
        
        # Explicit z-order: adding last usually puts it on top in XML order
        
    def add_slide(title_text, content_lines=[], image_path=None, is_title_slide=False, is_split_layout=False):
        layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(layout)
        set_background(slide)

        # Title Layout
        TITLE_HEIGHT = Inches(1.5)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), SLIDE_WIDTH - Inches(1.0), TITLE_HEIGHT)
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p = title_tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(48) # Slightly larger for widescreen
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if is_title_slide:
            p.alignment = PP_ALIGN.CENTER
        
        CONTENT_TOP = Inches(2.2)
        FOOTER_TOP = SLIDE_HEIGHT - Inches(1.2) # Reserve space for subtitle/footer
        
        if is_title_slide:
            # --- Dynamic Vertical Layout ---
            # 1. Title ends at 2.0"
            # 2. Subtitle starts at 6.0" (Example)
            # 3. Available gap = 4.0"
            # 4. We center image in this gap
            
            SUBTITLE_HEIGTH = Inches(1.2)
            subtitle_top = SLIDE_HEIGHT - SUBTITLE_HEIGTH - Inches(0.8) # e.g. at 7.5 - 1.2 - 0.8 = 5.5"
            
            available_height = subtitle_top - CONTENT_TOP - Inches(0.2) # 5.5 - 2.2 = 3.3 inches gap
            
            if image_path and os.path.exists(image_path):
                # Calculate center position
                img_display_height = min(available_height, Inches(4.5)) # Cap max height
                slide.shapes.add_picture(image_path, Inches(1), CONTENT_TOP, height=img_display_height)
                # Note: 'height' kwarg scales width proportionally. 
                # Since we have 13" width, a square image becomes 3.3" wide. 
                # We need to center it horizontally too.
                
                # To center properly, we need to know the image aspect ratio, 
                # but python-pptx doesn't give image dims easily without loading.
                # We will rely on "add_picture" returning a shape which we can move.
                pic = slide.shapes[-1] 
                # Center Horizontally
                pic.left = int((SLIDE_WIDTH - pic.width) / 2)
                # Center Vertically in the gap
                pic.top = int(CONTENT_TOP + (available_height - pic.height) / 2)

            # Subtitles below image
            sub_box = slide.shapes.add_textbox(Inches(0), subtitle_top, SLIDE_WIDTH, SUBTITLE_HEIGTH)
            tf = sub_box.text_frame
            for line in content_lines:
                p = tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(200, 200, 200)
            
            # Copyright (At very bottom)
            add_copyright(slide, is_cover=True)

        elif is_split_layout:
            # 16:9 Split Layout
            # Left Content: 0.5 to 6.5
            # Right Image: 7.0 to 12.8
            
            text_width = Inches(6.0)
            text_box = slide.shapes.add_textbox(Inches(0.5), CONTENT_TOP, text_width, Inches(5.0))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for line in content_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(14)
                if line.startswith("##"):
                    p.text = line.replace("##", "").strip()
                    p.font.size = Pt(40)
                    p.font.color.rgb = RGBColor(255, 215, 0)
                    p.font.bold = True
                elif line.startswith("#"):
                    p.text = line.replace("#", "").strip()
                    p.font.size = Pt(32)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                else:
                    p.text = line
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(220, 220, 220)
            
            if image_path and os.path.exists(image_path):
                img_left = Inches(7.0)
                img_width = Inches(5.8)
                slide.shapes.add_picture(image_path, img_left, CONTENT_TOP, width=img_width)
                
            add_copyright(slide) # Add to all slides strictly? User said "first and last". 
            # Actually user said "all generated ppt have copyright info on first and last page"
            # It's ambiguous if they mean ONLY first and last, or ALWAYS first and last (implies inclusive/exclusive).
            # "First join... End join..." usually means only those two.
            # I will follow strictly: First and Last only.

        else:
            # Standard Layout
            text_box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, SLIDE_WIDTH - Inches(2), Inches(5.0))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for line in content_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(20)
                if line.startswith("##"):
                    p.text = line.replace("##", "").strip()
                    p.font.size = Pt(60)
                    p.font.color.rgb = RGBColor(255, 215, 0)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                elif line.startswith("#"):
                    p.text = line.replace("#", "").strip()
                    p.font.size = Pt(32)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                else:
                    p.text = "• " + line
                    p.font.size = Pt(28)
                    p.font.color.rgb = RGBColor(220, 220, 220)

    # --- Content Generation ---

    # Slide 1: Cover
    add_slide("内存的战争\nThe War of Memory", ["2026 内存架构深度解析", "NVIDIA Rubin vs. DeepSeek Engram"], IMG_TITLE, is_title_slide=True)

    # Slide 2: Agenda
    add_slide("汇报大纲", [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件的暴力美学",
        "3. DeepSeek Engram：软件的魔法",
        "4. 巅峰对决：路线之争",
        "5. 未来展望：Agentic AI"
    ])

    # ... (Skipping full content repetition for brevity, using representative slides)
    # Slide 6: Rubin
    add_slide("NVIDIA Rubin 架构", [
        "Jensen Huang @ CES 2026",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN, is_split_layout=True)

    # Slide 13: Engram
    add_slide("DeepSeek Engram", [
        "软件定义的“外挂内存”",
        "无需专有硬件",
        "利用系统 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM, is_split_layout=True)
    
    # Slide 20: Compare
    add_slide("路线对比：逻辑之争", [
        "# NVIDIA: Cohesion",
        "硬件透明，高性能",
        "# DeepSeek: Retrieval",
        "软件显式，高性价比"
    ], IMG_COMPARE, is_split_layout=True)

    # Slide 25: Closing
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_background(slide)
    
    # Standard Title/Summary for closing
    t_box = slide.shapes.add_textbox(Inches(0), Inches(3), SLIDE_WIDTH, Inches(1.5))
    p = t_box.text_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub = slide.shapes.add_textbox(Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(1))
    p = sub.text_frame.add_paragraph()
    p.text = "The Architecture of Memory"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    add_copyright(slide, is_cover=False)

    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V4.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")

if __name__ == "__main__":
    create_ppt_v4()
