"""
Memory Architecture Report - V7 (40 Pages with Speaker Notes)
Generated following leijunskill standards:
- 40 pages with Lei Jun style speaker notes in PPT notes section
- Technical concept visualizations (§8)
- Signature only on Cover and Closing pages (§5)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"
IMG_MEMORY_WALL = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_memory_wall_1768530473358.png"
IMG_KV_CACHE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_kv_cache_1768530494022.png"
IMG_NGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_ngram_hashing_1768530510183.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_prefetch_strategy_1768530526058.png"

# --- CONSTANTS ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(1.5)
CONTENT_TOP = Inches(2.0)
CONTENT_HEIGHT = Inches(4.0)
SPLIT_TEXT_LEFT = Inches(0.5)
SPLIT_TEXT_WIDTH = Inches(5.5)
SPLIT_IMG_LEFT = Inches(6.5)
SPLIT_IMG_WIDTH = Inches(6.3)


def create_ppt_v7():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_notes(slide, notes_text):
        """Add speaker notes to slide - leijunskill §1.5"""
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes_text

    def add_signature(slide, is_cover=True):
        sig_width = Inches(5)
        sig_height = Inches(0.9)
        sig_left = (SLIDE_WIDTH - sig_width) / 2
        sig_top = SLIDE_HEIGHT - sig_height - Inches(0.15)
        box = slide.shapes.add_textbox(sig_left, sig_top, sig_width, sig_height)
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(0, 191, 255)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(255, 215, 0)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def add_title(slide, text, center=False):
        box = slide.shapes.add_textbox(Inches(0.5), TITLE_TOP, SLIDE_WIDTH - Inches(1), TITLE_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if center:
            p.alignment = PP_ALIGN.CENTER

    def add_content(slide, lines):
        box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, SLIDE_WIDTH - Inches(2), CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(16)
            if line.startswith("##"):
                p.text = line[2:].strip()
                p.font.size = Pt(54)
                p.font.color.rgb = RGBColor(255, 215, 0)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = "• " + line
                p.font.size = Pt(26)
                p.font.color.rgb = RGBColor(220, 220, 220)

    def add_split(slide, lines, img_path):
        box = slide.shapes.add_textbox(SPLIT_TEXT_LEFT, CONTENT_TOP, SPLIT_TEXT_WIDTH, CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(12)
            if line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(220, 220, 220)
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, SPLIT_IMG_LEFT, CONTENT_TOP, width=SPLIT_IMG_WIDTH)

    # ========== PAGE GENERATION (40 pages with notes) ==========

    # --- P1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "内存的战争\nThe Architecture of Memory", center=True)
    if os.path.exists(IMG_TITLE):
        s.shapes.add_picture(IMG_TITLE, Inches(4), Inches(2.2), height=Inches(3.0))
    sub = s.shapes.add_textbox(Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "NVIDIA Rubin vs. DeepSeek Engram | 2026 内存架构深度解析"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)
    add_notes(s, "朋友们，大家好！今天，我想和大家聊一个困扰AI行业多年的难题——内存。在2026年的今天，我们终于迎来了两个截然不同却殊途同归的解决方案。这，就是我们今天要讨论的内存的战争。")

    # --- P2: Agenda ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "汇报大纲")
    add_content(s, [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件霸权",
        "3. DeepSeek Engram：软件革命",
        "4. 路线对比：硬件 vs 软件",
        "5. 性能与经济分析",
        "6. 未来展望：Agentic AI"
    ])
    add_notes(s, "我们今天的分享分为以下几个部分。首先，我会介绍2026年AI行业面临的两大挑战——内存墙和上下文墙。接着，我们将深入分析NVIDIA的Rubin架构和DeepSeek的Engram技术。然后，我会对比这两条路线的优劣。最后，我们展望一下未来的Agentic AI时代。让我们一起来看。")

    # --- P3: Paradigm Shift ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2026：范式转移")
    add_content(s, [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ])
    add_notes(s, "朋友们，2026年，AI行业发生了一个根本性的转变。过去我们追求的是算力，是FLOPS。但现在，当模型参数突破万亿大关，我们发现——内存，才是新的瓶颈。记忆，成为了这个时代的石油。谁掌握了内存，谁就掌握了AI的未来。")

    # --- P4: Memory Wall ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Memory Wall：内存墙")
    add_split(s, [
        "#GPU 算力远超内存带宽",
        "数据传输成为瓶颈",
        "权重访问速度限制推理",
        "#算力越强，内存墙越高"
    ], IMG_MEMORY_WALL)
    add_notes(s, "我们来看这张图。左边是GPU，算力强大，金光闪闪。右边是内存，数据的仓库。但中间这堵墙——这就是内存墙。GPU再强，数据过不来，也是白搭。这意味着什么？这意味着，算力越强，这堵墙就越高。这是一个令人绝望的困境。")

    # --- P5: Context Wall ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Context Wall：上下文墙")
    add_content(s, [
        "KV Cache 随序列长度指数膨胀",
        "模型无法记住完整对话",
        "处理百万 Token 需天价计算",
        "## 失忆，是 AI 最大的痛点"
    ])
    add_notes(s, "除了内存墙，还有另一堵墙——上下文墙。当你和AI聊天的时候，每多说一句话，它需要记住的东西就指数级增长。结果呢？AI会忘记你之前说过的话。失忆，朋友们，这是AI最大的痛点。我们必须解决它。")

    # --- P6: Two Camps ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "两大阵营的回答")
    add_content(s, [
        "#NVIDIA (硬件派)",
        "Rubin 架构：极致协同，专有互联",
        "#DeepSeek (软件派)",
        "Engram 技术：算法优化，软件定义内存"
    ])
    add_notes(s, "面对这两堵墙，行业给出了两个方案。一边是NVIDIA，他们说：用最强的硬件，造最大的池子，用钞能力解决问题。另一边是DeepSeek，他们说：我们没有那么多钱，但我们有智慧，用软件来弥补硬件的不足。这，就是两大阵营的回答。")

    # --- P7: Rubin Architecture ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVIDIA Rubin 架构")
    add_split(s, [
        "#Jensen Huang @ CES 2026",
        "年度节奏发布",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN)
    add_notes(s, "让我们先看NVIDIA的方案。在2026年的CES上，老黄发布了Rubin架构。这不是一个简单的GPU升级，而是一整套芯片生态系统。他们把整个数据中心当作一块芯片来设计。这，就是极致的协同。")

    # --- P8: Six Chip Ecosystem ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "六大芯片生态系统")
    add_content(s, [
        "Rubin GPU - 核心计算引擎",
        "Vera CPU - 内存扩展",
        "NVLink 6 Switch - 高速互联",
        "ConnectX-9 SuperNIC - 网络加速",
        "BlueField-4 DPU - 上下文管理",
        "Spectrum-6 Ethernet Switch"
    ])
    add_notes(s, "这六大芯片，各司其职。Rubin GPU负责计算，Vera CPU负责内存扩展，NVLink 6负责芯片之间的高速互联。每一块芯片都是专门设计的，只为一个目标——让AI跑得更快，记得更多。")

    # --- P9: Rubin GPU ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin GPU：性能怪兽")
    add_content(s, [
        "## 3360 亿",
        "晶体管数量 (台积电 3nm)",
        "## 50 PFLOPS",
        "FP4 推理性能"
    ])
    add_notes(s, "大家看这个数字——3360亿！这是Rubin GPU的晶体管数量。采用台积电3nm工艺。50 PFLOPS的推理性能。是的，你没有听错，50 PFLOPS！这意味着什么？这意味着单块GPU的算力，相当于5年前一整个超算中心。")

    # --- P10: HBM4 Bandwidth ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：速度的革命")
    add_content(s, [
        "## 22 TB/s",
        "内存带宽 (vs Blackwell 8 TB/s)",
        "## 2.75x",
        "代际提升"
    ])
    add_notes(s, "再看这个数字——22 TB/s！这是HBM4的带宽。相比上一代Blackwell的8 TB/s，提升了2.75倍！朋友们，这不是小修小补，这是代际飞跃。数据像洪水一样在芯片里奔流，没有任何阻碍。")

    # --- P11: HBM4 Capacity ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：容量震撼")
    add_content(s, [
        "## 288 GB",
        "单卡 HBM4 显存容量",
        "可容纳 1440 亿参数模型权重",
        "MoE 路由决策可在微秒内完成"
    ])
    add_notes(s, "288 GB！单卡288 GB的显存。这意味着一块卡就能装下一个1440亿参数的模型。过去需要一整个机柜才能做的事情，现在一块卡就能搞定。这，就是技术的进步。")

    # --- P12: Vera CPU ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Vera CPU：AI 工厂专属")
    add_content(s, [
        "88 颗 Olympus 定制 ARM 核心",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "首款支持 FP8 的 CPU"
    ])
    add_notes(s, "但288 GB还不够。所以NVIDIA设计了Vera CPU，88颗定制ARM核心，配备1.5 TB的内存池。注意，是1.5 TB！这是GPU显存的5倍多。而且是首款支持FP8的CPU，专门为AI推理设计。")

    # --- P13: NVLink-C2C ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVLink-C2C：高速一致性")
    add_content(s, [
        "## 1.8 TB/s",
        "芯片间互联带宽",
        "相比 Grace-Blackwell 提升 2x",
        "GPU 可像访问本地内存一样访问 CPU 内存"
    ])
    add_notes(s, "1.8 TB/s的芯片间互联带宽！这意味着GPU访问CPU上的1.5 TB内存，就像访问自己的本地内存一样快。过去是两个孤岛，现在是一体化的大陆。这，就是NVLink-C2C的魔力。")

    # --- P14: NVL72 ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVL72：超级机柜")
    add_content(s, [
        "72 块 Rubin GPU + 36 块 Vera CPU",
        "## 54 TB",
        "机柜总内存容量",
        "统一内存寻址，透明访问"
    ])
    add_notes(s, "把这些芯片组合起来，就是NVL72超级机柜。72块GPU，36块CPU，总共54 TB的内存！统一寻址，透明访问。整个机柜就像一块超级大芯片。这，就是NVIDIA的暴力美学。")

    # --- P15: BlueField-4 ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "BlueField-4 与 ICMS")
    add_split(s, [
        "#G3.5 上下文存储层",
        "介于 HBM/DRAM 与 SSD 之间",
        "每 GPU 最高 16 TB 上下文",
        "KV Cache 成为共享资源"
    ], IMG_KV_CACHE)
    add_notes(s, "我们来看这张图，这是KV Cache的可视化。BlueField-4 DPU创造了一个新的存储层级——G3.5。每个GPU可以访问高达16 TB的上下文存储。KV Cache不再是本地资源，而是共享的基础设施。这，彻底解决了长期记忆的问题。")

    # --- P16: Rubin Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin 硬件总结")
    add_content(s, [
        "#核心理念：物理统一内存空间",
        "22 TB/s HBM4 + 1.8 TB/s C2C + 16 TB ICMS",
        "#代价：全栈锁定 NVIDIA 生态",
        "高 CapEx，高性能"
    ])
    add_notes(s, "总结一下NVIDIA的方案：用最强的硬件，创造一个物理上统一的内存空间。22 TB/s的HBM4，1.8 TB/s的C2C互联，16 TB的上下文存储。代价是什么？全栈锁定NVIDIA生态，高额的资本投入。这是巨头的游戏。")

    # --- P17: DeepSeek Engram ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "DeepSeek Engram")
    add_split(s, [
        "#软件定义的外挂内存",
        "无需专有硬件",
        "利用主机 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM)
    add_notes(s, "现在，让我们看看另一条路。DeepSeek的Engram技术。他们说：我们没有NVIDIA那么多钱，但我们有聪明的脑袋。我们用软件，把主机上便宜的DDR内存，变成GPU的外挂记忆。不需要专有硬件，普通的电脑就能用。")

    # --- P18: Core Insight ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "核心洞察")
    add_content(s, [
        "LLM 大量计算浪费在重构静态知识",
        "为什么不直接查找?",
        "## Reconstruct vs. Lookup",
        "用查表替代计算"
    ])
    add_notes(s, "Engram的核心洞察是什么？DeepSeek发现，大语言模型90%的计算，都在干同一件事——重构那些早就知道的知识。这是浪费！为什么不直接查表呢？Reconstruct vs Lookup，重构 vs 查找。用查表替代计算，这，就是Engram的精髓。")

    # --- P19: N-gram Hashing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "N-gram Hashing 原理")
    add_split(s, [
        "#输入 Token 序列编码为 Hash",
        "确定性、上下文相关的 ID",
        "用 Hash 作为 Key 查询",
        "## O(1) 复杂度查表"
    ], IMG_NGRAM)
    add_notes(s, "我们来看这张图，这是N-gram Hashing的流程。输入的Token序列，通过Hash函数，变成一个确定性的ID。用这个ID去查表，取出对应的知识向量。关键是什么？O(1)复杂度！无论表有多大，查询时间都是恒定的。这，就是算法的力量。")

    # --- P20: Embedding Table ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Embedding Table：知识仓库")
    add_content(s, [
        "## 100B 参数",
        "存放在主机 DDR 内存中",
        "27B 活跃模型 + 100B 知识表",
        "小模型拥有万亿知识"
    ])
    add_notes(s, "大家看这个数字——100B参数的知识表，存放在主机的DDR内存中。一个27B的小模型，加上100B的知识表，总共127B参数。效果呢？堪比万亿参数的大模型！小身材，大智慧。这，就是Engram的魔法。")

    # --- P21: PCIe Bottleneck ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "克服 PCIe 瓶颈")
    add_split(s, [
        "#PCIe 带宽远低于 NVLink",
        "解决方案：确定性预取",
        "提前计算 Hash ID",
        "异步获取数据"
    ], IMG_PREFETCH)
    add_notes(s, "我们来看这张图。上面是GPU计算的进度条，下面是内存预取的进度条。问题是什么？PCIe比NVLink慢得多。解决方案？确定性预取！因为我们知道接下来需要哪些数据，所以可以提前获取。当GPU需要的时候，数据已经准备好了。完美同步！")

    # --- P22: Prefetch Effect ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "预取策略效果")
    add_content(s, [
        "在 GPU 执行早期层时",
        "异步从 DDR 获取后续所需数据",
        "## 吞吐量损失 < 3%",
        "廉价 DDR 成为 HBM 的可行扩展"
    ])
    add_notes(s, "效果如何？吞吐量损失不到3%！是的，你没有听错，只有3%。这意味着什么？这意味着便宜的DDR内存，现在可以作为昂贵的HBM的可行扩展。软件的智慧，弥补了硬件的不足。")

    # --- P23: Engram vs RAG ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram vs RAG")
    add_content(s, [
        "#传统 RAG：文档级检索",
        "在 Transformer 外部操作",
        "#Engram：Token 级查表",
        "在 Transformer 内部融合",
        "粒度更细，延迟更低"
    ])
    add_notes(s, "有人会问：这和RAG有什么区别？区别很大。传统RAG是文档级别的检索，在模型外部操作。而Engram是Token级别的查表，直接嵌入到Transformer内部。粒度更细，延迟更低，效果更好。")

    # --- P24: Engram Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram 软件总结")
    add_content(s, [
        "#核心理念：算法替代硬件",
        "N-gram Hashing + 预取重叠",
        "#代价：需修改模型架构",
        "低成本，高可及性"
    ])
    add_notes(s, "总结一下DeepSeek的方案：用算法替代硬件，用智慧弥补资源的不足。代价是什么？需要修改模型架构。但好处是：低成本，高可及性。普通人也能用得起的AI，这，就是Engram的意义。")

    # --- P25: Route Comparison ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "路线对比")
    add_split(s, [
        "#NVIDIA: Cohesion (内聚)",
        "硬件透明，模型无感知",
        "#DeepSeek: Retrieval (检索)",
        "软件显式，模型需适配"
    ], IMG_COMPARE)
    add_notes(s, "现在，让我们对比一下这两条路线。NVIDIA的方案是内聚，硬件层面解决问题，模型不需要知道内存在哪里。DeepSeek的方案是检索，软件层面解决问题，模型需要显式地调用查表操作。两种哲学，两条道路。")

    # --- P26: Memory Type Comparison ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "目标记忆类型对比")
    add_content(s, [
        "#Rubin: 动态 KV Cache",
        "对话上下文 (短期记忆)",
        "每次 Token 生成都需重读",
        "#Engram: 静态 Knowledge",
        "世界知识 (长期记忆)",
        "训练时固化，推理时查表"
    ])
    add_notes(s, "更重要的是，这两个方案解决的是不同类型的记忆问题。Rubin解决的是KV Cache，就是对话的短期记忆。Engram解决的是静态知识，就是世界的长期记忆。两者其实是互补的！")

    # --- P27: Tech Specs Comparison ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "技术参数对比")
    add_content(s, [
        "#NVIDIA Rubin",
        "22 TB/s 带宽, 54 TB 容量, 专有硬件",
        "#DeepSeek Engram",
        "PCIe 带宽, 无限 DDR, 开源软件"
    ])
    add_notes(s, "从技术参数来看：Rubin有22 TB/s的带宽，54 TB的容量，但需要专有硬件。Engram只有PCIe的带宽，但可以使用无限量的DDR内存，而且是开源软件。性能 vs 可及性，这是核心的权衡。")

    # --- P28: Use Case Comparison ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "适用场景对比")
    add_content(s, [
        "#Rubin 适合",
        "超大规模云厂商 (Microsoft/AWS)",
        "需要极致性能和 SLA 保障",
        "#Engram 适合",
        "资源受限环境",
        "消费级硬件部署"
    ])
    add_notes(s, "那么，谁应该选择哪个方案？Rubin适合超大规模云厂商，微软、亚马逊这些巨头。他们需要极致性能，需要SLA保障。Engram适合资源受限的环境，创业公司、研究机构、个人开发者。各取所需。")

    # --- P29: Complementary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "互补而非替代")
    add_content(s, [
        "Rubin 解决 KV Cache 容量",
        "Engram 解决知识存储",
        "## 两者可以结合",
        "DDR 存知识，HBM 存上下文"
    ])
    add_notes(s, "朋友们，我想强调的是：这两个方案不是非此即彼的关系，而是可以结合的！用Engram把知识存到DDR里，用Rubin的HBM存对话上下文。各司其职，物尽其用。这，才是最优解。")

    # --- P30: Open vs Closed ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "开源 vs 闭源")
    add_content(s, [
        "#NVIDIA",
        "专有生态，高锁定，高壁垒",
        "#DeepSeek",
        "开源权重，低门槛，普惠创新"
    ])
    add_notes(s, "最后，从生态的角度来看。NVIDIA是专有生态，一旦选择，就很难离开。DeepSeek是开源的，所有人都可以用，都可以改进。一边是高壁垒，一边是普惠创新。你会选择哪一边？")

    # --- P31: Engram Benchmark ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram 性能实测")
    add_content(s, [
        "## +5.0",
        "BBH (复杂推理)",
        "## +3.4",
        "MMLU (知识问答)"
    ])
    add_notes(s, "让我们看看实测数据。在BBH复杂推理测试中，Engram比同等规模的MoE模型高出5.0分！在MMLU知识问答测试中，高出3.4分！一个27B的小模型，打败了那些庞然大物。这，就是算法的力量。")

    # --- P32: Rubin Economics ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin 经济效益")
    add_content(s, [
        "## 10x",
        "Token 成本下降 (vs Blackwell)",
        "## $5B/$100M",
        "投资回报率 (Token 收入/基础设施投入)"
    ])
    add_notes(s, "再看Rubin的经济效益。相比Blackwell，Token成本下降了10倍！投资回报率是多少？每投入1亿美元基础设施，可以产生50亿美元的Token收入。这是一个50倍的ROI！对于大厂来说，这是不能拒绝的买卖。")

    # --- P33: AI Democratization ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "AI 民主化")
    add_content(s, [
        "无需 H100/Rubin",
        "RTX 4090 集群跑 GPT-5 级模型",
        "## 打破 NVIDIA 税"
    ])
    add_notes(s, "但对于我们普通人呢？Engram带来了AI的民主化。你不需要H100，不需要Rubin。用RTX 4090的集群，就能跑GPT-5级别的模型！打破NVIDIA税，让每个人都能用得起最强的AI。这，才是真正的普惠。")

    # --- P34: Geopolitics ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "地缘政治意义")
    add_content(s, [
        "面对出口管制",
        "Engram 是突围方案",
        "DeepSeek V3 仅用 $5.58M 训练",
        "## More with Less"
    ])
    add_notes(s, "这里还有一个地缘政治的角度。在出口管制的背景下，很多地方拿不到最新的GPU。Engram提供了一个突围方案：用软件弥补硬件的不足。DeepSeek V3仅用558万美元就完成了训练。More with Less，用更少做更多。这，是一种智慧。")

    # --- P35: Agentic AI ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "未来：Agentic AI")
    add_content(s, [
        "从 Chatbot 到 Agent",
        "多步推理，长期记忆",
        "管理复杂工作流",
        "## Context is the new Currency"
    ])
    add_notes(s, "最后，让我们展望一下未来。AI正在从Chatbot演进到Agent。Agent不只是回答问题，而是要完成复杂的任务，需要多步推理，需要长期记忆。在这个新时代，上下文就是新的货币。谁掌握了上下文，谁就掌握了未来。")

    # --- P36: Million Token Context ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "百万 Token 上下文")
    add_content(s, [
        "Rubin CPX 可摄入完整代码库",
        "一小时视频作为输入",
        "## 将整个仓库作为工作记忆"
    ])
    add_notes(s, "想象一下：AI可以一次性处理整个代码仓库，一小时的视频。将所有内容作为工作记忆。你可以问它任何问题，它都能基于完整的上下文来回答。这，就是百万Token上下文带来的可能性。")

    # --- P37: Shared Memory Pool ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "共享记忆池")
    add_content(s, [
        "BlueField-4 管理 KV Cache",
        "多 Agent 共享，无需重算",
        "## 协作式 AI 成为可能"
    ])
    add_notes(s, "更进一步，多个Agent可以共享同一个记忆池。一个Agent生成的KV Cache，另一个Agent可以直接使用，无需重新计算。这意味着什么？协作式AI成为可能。多个Agent像一个团队一样协同工作。")

    # --- P38: Long-term Consistency ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "长期一致性记忆")
    add_content(s, [
        "Engram 支持三种记忆类型",
        "Episodic / Semantic / Procedural",
        "## +15 分 (仅用 1% Token)"
    ])
    add_notes(s, "Engram还支持三种类型的记忆：情景记忆、语义记忆、程序记忆。在记忆一致性测试中，比全上下文基线高出15分，却只用了1%的Token！这，就是高效记忆的威力。")

    # --- P39: Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "总结：内存革命")
    add_content(s, [
        "Rubin 定义了性能上限",
        "Engram 拉高了可及下限",
        "## 殊途同归：解决 AI 失忆"
    ])
    add_notes(s, "朋友们，让我来总结一下。NVIDIA的Rubin定义了性能的天花板，DeepSeek的Engram拉高了可及性的地板。两条不同的路，却指向同一个目标——解决AI的失忆问题。这，就是2026年内存革命的意义。")

    # --- P40: Closing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
    p = box.text_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "The Architecture of Memory"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=False)
    add_notes(s, "朋友们，这就是我们今天分享的全部内容。无论你选择硬件的霸道，还是软件的智慧，这都是一个记忆永驻的新时代。感谢大家的聆听，谢谢！")

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V7_with_notes.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_ppt_v7()
