
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Image Paths from previous generation
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"

def create_ppt_v2():
    prs = Presentation()
    
    # --- Skill Rule: Defensive Style Enforcement ---
    # We will set background manually for every slide to ensure dark mode
    BG_COLOR = RGBColor(15, 15, 25) # Deep LeiJun Blue/Black

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_slide(title_text, content_lines=[], image_path=None, is_title_slide=False, is_split_layout=False):
        if is_title_slide:
            layout = prs.slide_layouts[6] # Blank for total control
        else:
            layout = prs.slide_layouts[6] # Using Blank for all slots to enforce "Absolute Coordinates" rule

        slide = prs.slides.add_slide(layout)
        set_background(slide)

        # --- Skill Rule: Title Zone (Safe Zone: Top 0 - 2.0 inches) ---
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p = title_tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        # --- Skill Rule: Style Enforcement (Explicit White) ---
        p.font.color.rgb = RGBColor(255, 255, 255)
        # Center title for visual impact
        if is_title_slide:
            p.alignment = PP_ALIGN.CENTER
        
        
        # --- Skill Rule: Content Zone (Start > 2.2 inches) ---
        CONTENT_TOP = Inches(2.5) 
        
        if is_title_slide:
            # Special Handling for Cover
            if image_path and os.path.exists(image_path):
                # Hero Image Centered
                left = Inches(1)
                top = Inches(2.2)
                width = Inches(8)
                slide.shapes.add_picture(image_path, left, top, width=width)
                
            # Subtitles below image
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(1))
            tf = sub_box.text_frame
            for line in content_lines:
                p = tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(200, 200, 200)

        elif is_split_layout:
            # --- Skill Rule: Split Layout & Safe Zones ---
            # Text uses left 5 inches
            # Image uses right 3.5 inches
            # Gap of > 0.5 inches
            
            # 1. Text Content (Left)
            text_box = slide.shapes.add_textbox(Inches(0.5), CONTENT_TOP, Inches(4.5), Inches(4.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for line in content_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(14) # Spacing
                
                if line.startswith("##"):
                    p.text = line.replace("##", "").strip()
                    p.font.size = Pt(36)
                    p.font.color.rgb = RGBColor(255, 215, 0) # Gold
                    p.font.bold = True
                elif line.startswith("#"):
                    p.text = line.replace("#", "").strip()
                    p.font.size = Pt(28)
                    p.font.color.rgb = RGBColor(255, 255, 255) # White
                    p.font.bold = True
                else:
                    p.text = line
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(220, 220, 220) # Light Grey
            
            # 2. Image Content (Right - Absolute Coordinates)
            if image_path and os.path.exists(image_path):
                img_left = Inches(5.5) # 0.5 inch gap from 5.0
                img_width = Inches(4.0)
                slide.shapes.add_picture(image_path, img_left, CONTENT_TOP, width=img_width)

        else:
            # Standard List Layout (Centered or Full Width)
            # Use a slightly narrower box to avoid edge crowding
            text_box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, Inches(8), Inches(4.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for line in content_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(20)
                
                if line.startswith("##"):
                    p.text = line.replace("##", "").strip()
                    p.font.size = Pt(54) # Big numbers
                    p.font.color.rgb = RGBColor(255, 215, 0)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                elif line.startswith("#"):
                    p.text = line.replace("#", "").strip()
                    p.font.size = Pt(32)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                else:
                    p.text = "• " + line # Manual bullet
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(220, 220, 220)

    # --- Recreating Slides with V2 Architecture ---

    # Slide 1: Title
    add_slide("内存的战争\nThe War of Memory", ["2026 内存架构深度解析", "NVIDIA Rubin vs. DeepSeek Engram"], IMG_TITLE, is_title_slide=True)

    # Slide 2: Agenda
    add_slide("汇报大纲", [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件的暴力美学",
        "3. DeepSeek Engram：软件的魔法",
        "4. 巅峰对决：路线之争",
        "5. 未来展望：Agentic AI"
    ])

    # Slide 3: Intro
    add_slide("2026：范式转移", [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ])

    # Slide 4: The Walls (Split layout logic even without img for structure)
    add_slide("两堵高墙", [
        "# Memory Wall (内存墙)",
        "权重访问速度限制了推理速度",
        "# Context Wall (上下文墙)",
        "序列长度限制了思考深度",
        "我们需要：更快的读写，更大的容量"
    ])

    # Slide 5: The Two Players
    add_slide("两大阵营", [
        "# NVIDIA (硬件派)",
        "集成电路，专有互联",
        "# DeepSeek (软件派)",
        "算法优化，软件定义内存"
    ])

    # Slide 6: NVIDIA Rubin (Split Layout with Image)
    add_slide("NVIDIA Rubin 架构", [
        "Jensen Huang @ CES 2026",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN, is_split_layout=True)

    # Slide 7: Rubin GPU Specs
    add_slide("Rubin GPU：性能怪兽", [
        "## 3360 亿",
        "晶体管数量 (Transistors)",
        "## 50 PFLOPS",
        "FP4 推理性能"
    ])

    # Slide 8: HBM4 Revolution
    add_slide("HBM4：速度的革命", [
        "## 22 TB/s",
        "内存带宽 (Bandwidth)",
        "## 288 GB",
        "单卡显存容量"
    ])

    # Slide 9: Vera CPU
    add_slide("Vera CPU：内存扩展", [
        "专为 AI Factory 设计",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "NVLink-C2C 互联 @ 1.8 TB/s"
    ])

    # Slide 11: The Cost
    add_slide("成本与代价", [
        "高性能 = 高投入",
        "专有硬件确立护城河",
        "CapEx (资本支出) 巨大"
    ])

    # Slide 13: DeepSeek Engram (Split Layout with Image)
    add_slide("DeepSeek Engram", [
        "软件定义的“外挂内存”",
        "无需专有硬件",
        "利用系统 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM, is_split_layout=True)

    # Slide 14: Core Insight
    add_slide("核心洞察", [
        "LLM 90% 的计算在“重构知识”",
        "为什么不直接“查找”？",
        "# Reconstruct vs. Lookup"
    ])

    # Slide 15: Implementation
    add_slide("实现原理：Engram", [
        "N-gram Hashing (N元语法哈希)",
        "Embedding Table 存放在 DDR",
        "Context-Aware Gating"
    ])

    # Slide 16: Prefetching Strategy
    add_slide("克服 PCIe 瓶颈", [
        "PCIe 比 NVLink 慢得多",
        "解决方案：",
        "## Deterministic Prefetching",
        "异步预取，吞吐量损失 < 3%"
    ])

    # Slide 17: Performance
    add_slide("性能实测：以小博大", [
        "## +5.0",
        "BBH (复杂推理)",
        "## +3.4",
        "MMLU (知识问答)"
    ])

    # Slide 19: Economic Impact
    add_slide("经济学：AI 民主化", [
        "无需 H100/Rubin",
        "消费级显卡 + 大内存",
        "RTX 4090 集群跑 GPT-5"
    ])

    # Slide 20: Strategic Comparison (Split Layout)
    add_slide("路线对比：逻辑之争", [
        "# NVIDIA: Cohesion",
        "硬件透明，高性能",
        "# DeepSeek: Retrieval",
        "软件显式，高性价比"
    ], IMG_COMPARE, is_split_layout=True)

    # Slide 24: Conclusion
    add_slide("总结：内存革命", [
        "Rubin 定义了上限",
        "Engram 拉高了下限",
        "## 遗忘的时代结束了"
    ])

    # Slide 25: Closing
    add_slide("谢谢", [
        "The Architecture of Memory",
        "## 谢谢大家！"
    ], IMG_TITLE, is_title_slide=True)

    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V2.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")

if __name__ == "__main__":
    create_ppt_v2()
