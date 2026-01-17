"""
Memory Architecture Report - V9 (40 Pages)
Generated with leijunskill v2.0 enhancements:
- §1.5 Persona Building (真诚的奋斗者)
- §1.6 Lei's Comparison Method (雷氏对比法)
- §1.7 Peak-End Rule (峰终定律)
- §1.8 Emotional Resonance (情感共鸣)
- §2.1 Deletion Rule (删除法则)
- §2.4 Font Size Contrast (字号对比法则)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"
IMG_MEMORY_WALL = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_memory_wall_1768530473358.png"
IMG_KV_CACHE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_kv_cache_1768530494022.png"
IMG_NGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_ngram_hashing_1768530510183.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_prefetch_strategy_1768530526058.png"

# --- CONSTANTS (leijunskill §6) ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)

# --- FONT SIZE CONTRAST (leijunskill §2.4) ---
TITLE_SIZE = Pt(54)
BIG_NUMBER_SIZE = Pt(96)  # Core numbers 80-120pt
SUBTITLE_SIZE = Pt(32)
BODY_SIZE = Pt(26)
CAPTION_SIZE = Pt(16)

# Colors
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(255, 215, 0)
LIGHT_GRAY = RGBColor(220, 220, 220)
TIANYI_BLUE = RGBColor(0, 191, 255)


def create_ppt_v9():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    def add_signature(slide, is_cover=True):
        """Signature only on cover and closing - leijunskill §5"""
        box = slide.shapes.add_textbox(
            (SLIDE_WIDTH - Inches(5)) / 2,
            SLIDE_HEIGHT - Inches(1),
            Inches(5), Inches(0.9)
        )
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = TIANYI_BLUE
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(24)
        p2.font.color.rgb = GOLD
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def title_slide(title, subtitle, img_path, notes):
        """Cover slide with centered layout"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), SLIDE_WIDTH - Inches(1), Inches(1.8))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Image
        if img_path and os.path.exists(img_path):
            s.shapes.add_picture(img_path, Inches(4), Inches(2.3), height=Inches(2.8))
        # Subtitle
        sub = s.shapes.add_textbox(Inches(0), Inches(5.4), SLIDE_WIDTH, Inches(0.8))
        p = sub.text_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_signature(s, is_cover=True)
        add_notes(s, notes)
        return s

    def content_slide(title, bullets, notes):
        """Standard content slide with left-aligned bullets"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Content
        content_box = s.shapes.add_textbox(Inches(1), Inches(2.0), SLIDE_WIDTH - Inches(2), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(14)
            if line.startswith("##"):
                para.text = line[2:].strip()
                para.font.size = BIG_NUMBER_SIZE
                para.font.color.rgb = GOLD
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = SUBTITLE_SIZE
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = "• " + line
                para.font.size = BODY_SIZE
                para.font.color.rgb = LIGHT_GRAY
        add_notes(s, notes)
        return s

    def data_slide(title, number, unit, explanation, notes):
        """Big number slide - leijunskill §2.4 Font Size Contrast"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1.2))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Big Number (Core: 80-120pt Gold)
        num_box = s.shapes.add_textbox(Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(2))
        p = num_box.text_frame.add_paragraph()
        p.text = number
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER
        # Unit
        unit_box = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
        p = unit_box.text_frame.add_paragraph()
        p.text = unit
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        # Explanation
        exp_box = s.shapes.add_textbox(Inches(1), Inches(5.2), SLIDE_WIDTH - Inches(2), Inches(1))
        p = exp_box.text_frame.add_paragraph()
        p.text = explanation
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_notes(s, notes)
        return s

    def split_slide(title, bullets, img_path, notes):
        """Left text, right image - leijunskill §6.3"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        # Title
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = WHITE
        # Left content
        left_box = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for line in bullets:
            para = tf.add_paragraph()
            para.space_after = Pt(12)
            if line.startswith("#"):
                para.text = line[1:].strip()
                para.font.size = Pt(28)
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                para.text = line
                para.font.size = Pt(22)
                para.font.color.rgb = LIGHT_GRAY
        # Right image
        if img_path and os.path.exists(img_path):
            s.shapes.add_picture(img_path, Inches(6.5), Inches(2.0), width=Inches(6.3))
        add_notes(s, notes)
        return s

    def closing_slide(title, subtitle, notes):
        """Closing slide with signature"""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        sub = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
        p = sub.text_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        add_signature(s, is_cover=False)
        add_notes(s, notes)
        return s

    # ==================== 40 PAGES GENERATION ====================

    # --- P1: Cover (峰终定律：开场要震撼) ---
    title_slide(
        "内存的战争\nThe Architecture of Memory",
        "NVIDIA Rubin vs. DeepSeek Engram | 2026 内存架构深度解析",
        IMG_TITLE,
        """朋友们，大家好！

今天，我想和大家讲一个关于"记忆"的故事。

我相信在座的每一位都遇到过这样的情况：你和AI聊了半个小时，聊得热火朝天。突然你问它："我刚才说的那个想法，你觉得怎么样？" 它愣了两秒，说："抱歉，你说过什么想法吗？"

这就是AI的失忆症。而今天，我们要聊的就是——如何治好它。

两大阵营，给出了两个截然相反的答案。一个用钱砸出一条路，一个用脑子想出一条路。这，就是内存的战争。"""
    )

    # --- P2: Agenda ---
    content_slide("汇报大纲", [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件的暴力美学",
        "3. DeepSeek Engram：软件的智慧",
        "4. 巅峰对决：路线之争",
        "5. 性能与经济：数据说话",
        "6. 未来展望：Agentic AI"
    ], """我们今天的分享，就像一场探险之旅。

首先，我带大家认识两堵高墙——内存墙和上下文墙。这是AI发展路上最大的拦路虎。

然后，我们看看两位英雄如何翻越这堵墙。一个是财大气粗的NVIDIA，另一个是聪明绝顶的DeepSeek。

接着，我们把两条路线放在一起PK，看看各自的优劣。

最后，我们一起展望未来——Agentic AI时代即将到来。

准备好了吗？Let's go！""")

    # --- P3: Paradigm Shift ---
    content_slide("2026：范式转移", [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ], """朋友们，2026年，AI行业发生了一件大事。

过去十年，我们一直在追逐一个数字——FLOPS，每秒浮点运算次数。谁的算力高，谁就是王者。

但现在，规则变了。就像100年前，石油取代煤炭成为工业的血液。今天，记忆正在取代算力，成为AI的新命脉。

为什么？因为模型参数已经突破万亿大关。算力再强，记不住东西，也是白搭。

这不是我一个人的判断，这是整个行业的共识。记忆，是新的石油。""")

    # --- P4: Memory Wall (图文页 + 生活化比喻) ---
    split_slide("Memory Wall：内存墙", [
        "#GPU 算力远超内存带宽",
        "数据传输成为瓶颈",
        "权重访问速度限制推理",
        "#算力越强，墙越高"
    ], IMG_MEMORY_WALL,
    """我们来看这张图。

想象一下早高峰的地铁站。左边是源源不断涌来的乘客——这就是GPU产生的算力。右边是站台——这就是内存里的数据。中间这道闸机，就是带宽。

问题来了：乘客太多，闸机太少。大家挤在那里，谁也过不去。

这就是内存墙。我们的GPU越来越强，就像乘客越来越多。但闸机数量没怎么变。结果呢？算力空转，数据过不来。

更可怕的是：你越努力提升算力，这堵墙就越高。这是一个让工程师们抓狂的死循环。""")

    # --- P5: Context Wall ---
    content_slide("Context Wall：上下文墙", [
        "KV Cache 随序列长度指数膨胀",
        "模型记不住完整对话",
        "处理百万 Token 需天价计算",
        "## 失忆，是 AI 最大的痛点"
    ], """第二堵墙更扎心——上下文墙。

我给大家讲一个真实的故事。我有个同事，用AI帮他写季度报告。写到一半，他让AI修改第一段。AI回答说："抱歉，我不记得第一段写了什么。"

他当场就崩溃了。

这就是上下文墙的问题。AI的"工作记忆"是有限的。你和它聊得越多，它需要记住的东西就越多。这些记忆叫KV Cache，会随着对话长度指数级膨胀。

就像你的手机，照片越存越多，内存越来越满。最后不得不删照片。AI也一样，聊着聊着，就把前面的忘了。

朋友们，失忆，这是AI最大的痛点。我们必须解决它。""")

    # --- P6: Two Camps (雷氏对比法) ---
    content_slide("两大阵营的回答", [
        "#NVIDIA (硬件派)",
        "Rubin 架构：极致协同，专有互联",
        "#DeepSeek (软件派)",
        "Engram 技术：算法优化，软件定义内存"
    ], """面对这两堵墙，行业出现了两个阵营。

（雷氏对比法：二元对立）

一边是NVIDIA，就像一个土豪。他们说："墙太高？没事，我有钱，我造电梯！用最贵的材料，打造最强的梯子。钱不是问题。"

另一边是DeepSeek，像一个聪明的穷学生。他们说："我们确实没那么多钱。但我们有脑子！找条隧道绕过去不行吗？"

一个靠钱，一个靠脑子。两条完全不同的路，目标却一样——翻过那堵墙。

说实话，在做这个分析的时候，我失眠了好几天。因为这两条路，代表着AI未来十年的走向。""")

    # --- P7: NVIDIA Rubin Architecture ---
    split_slide("NVIDIA Rubin 架构", [
        "#Jensen Huang @ CES 2026",
        "年度节奏发布",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN,
    """让我们先看土豪的方案。

2026年CES上，老黄穿着他标志性的皮衣走上台。他说了一句话，让全场沸腾："我们不是在做GPU，我们是在做一个超级大脑。"

Rubin架构的思路简单粗暴：既然问题是内存不够，那就把所有的芯片都连起来，共享内存。GPU、CPU、网卡、DPU……全部打通。

就像把一栋栋独立的房子，改造成一个巨型公寓楼。每个房间都能用整栋楼的水电气。

这叫"极致协同"。听起来简单，做起来要命。但NVIDIA做到了。""")

    # --- P8: Six Chip Ecosystem ---
    content_slide("六大芯片生态系统", [
        "Rubin GPU - 核心算力引擎",
        "Vera CPU - 内存扩展",
        "NVLink 6 Switch - 高速互联",
        "ConnectX-9 SuperNIC - 网络加速",
        "BlueField-4 DPU - 上下文管理",
        "Spectrum-6 Ethernet Switch"
    ], """这六块芯片，就像复仇者联盟。

Rubin GPU是钢铁侠，战斗力最强，负责核心计算。
Vera CPU是绿巨人，力大无穷，扛着1.5TB的内存。
NVLink 6是蜘蛛侠的蛛丝，把所有人连在一起。
ConnectX-9是鹰眼，眼观六路，网络通信全靠它。
BlueField-4是幻视，管理着所有人的记忆。

六个超级英雄，组成一支无敌战队。这不是我一个人的成绩，这是NVIDIA 3万名工程师，日夜奋战的结果。""")

    # --- P9: Rubin GPU - Big Number (峰终定律：高潮点) ---
    data_slide("Rubin GPU", "3360亿", "晶体管数量 (台积电 3nm)", "相比上一代提升 1.6x",
    """大家看这个数字——3360亿！

（数字具象化）
3360亿是什么概念？如果每个晶体管是一粒沙子，这些沙子可以填满40个奥运会游泳池！然后，全部塞进一块比你手掌还小的芯片里。

这不是科幻电影，这是2026年的现实。

我们的工程师告诉我，为了实现这个目标，他们连续攻关了18个月，试了超过500种设计方案。最终找到了最优解。

是的，你没有听错。3360亿晶体管，小米手机都装不下这么多照片。""")

    # --- P10: HBM4 Bandwidth ---
    data_slide("HBM4 带宽革命", "22 TB/s", "内存带宽", "vs Blackwell 8 TB/s，提升 2.75x",
    """再看这个数字——22 TB/s！

（数字具象化）
22 TB/s是什么概念？让我给大家算一笔账。

一部高清电影大概5GB。22 TB/s意味着：每秒钟可以传输4400部高清电影。一眨眼的功夫，整个Netflix的片库就传完了。

还记得刚才说的地铁闸机吗？以前是8个闸机，现在变成22个，还都是快速通道。乘客再多，也不用挤了。

这不是小改进，这是代际飞跃。从8到22，提升了2.75倍！

（工程师思维）
而且，这是通过硅片层面的优化实现的，不是靠压缩算法凑出来的。真材实料。""")

    # --- P11: HBM4 Capacity ---
    data_slide("HBM4 容量震撼", "288 GB", "单卡显存", "可容纳 1440 亿参数模型",
    """288 GB！单卡288 GB！

（对比故事）
我给大家讲个故事。三年前，我们团队要跑一个700亿参数的模型。买了8块A100，每块80GB，凑了640GB才勉强够用。机器塞满了一整个机柜，电费一个月好几万。

现在呢？

一块Rubin，288GB，装下1440亿参数，一块卡就够了。

三年前需要一个机柜，今天只要一块卡。这就是技术进步的速度。

朋友们，看到这个数字，我真的很感慨。我们这一代人，正在见证历史。""")

    # --- P12: Vera CPU ---
    data_slide("Vera CPU", "1.5 TB", "LPDDR5X 内存池", "88 颗 Olympus ARM 核心",
    """但288 GB还不够！

所以NVIDIA造了第二个武器——Vera CPU，配备1.5 TB的内存。

（生活化比喻）
如果GPU是一个超级大厨，手边的料理台（HBM）放不下太多食材。没关系，旁边还有一个巨大的冷库（Vera内存），什么都往里塞。需要的时候，通过高速传送带（NVLink）瞬间送过来。

关键是，这个冷库就在大厨隔壁，不是隔着几条街的超市。拿东西快得很。

88颗定制ARM核心，专门负责搬运数据。不做计算，就做一件事——快！""")

    # --- P13: NVLink-C2C ---
    data_slide("NVLink-C2C", "1.8 TB/s", "芯片间互联带宽", "相比 Grace-Blackwell 提升 2x",
    """1.8 TB/s的芯片间互联！

这是什么概念？从GPU到CPU的数据传输速度，和GPU访问自己内部显存一样快。

（生活化比喻）
就像两栋楼之间架了一座天桥，走过去只要10秒。以前呢？要下楼、过马路、等红灯、上楼，至少10分钟。

对于程序来说，它感觉不到内存在两个不同的芯片上。就像你住公寓，不需要知道自来水从哪个水厂来，打开水龙头就有。

这就是硬件工程师的魔法——把复杂的事情，变得简单透明。""")

    # --- P14: NVL72 Rack ---
    data_slide("NVL72 超级机柜", "54 TB", "总内存容量", "72 GPU + 36 CPU，统一寻址",
    """把这些加起来，就是NVL72超级机柜。

54 TB的内存！

（数字具象化）
假设一个人的大脑记忆容量是2.5 PB，那么NVL72相当于人类大脑的2%。听起来不多？

但别忘了：人脑花了几百万年进化，NVL72只用了几年开发。

（共同体意识）
更重要的是，这54 TB对所有GPU来说都是透明可见的。就像一个超大的共享网盘，谁需要就去取。

这不是某一个团队的成绩，是NVIDIA整个公司的协同作战。硬件、软件、网络，三军联合。""")

    # --- P15: BlueField-4 ---
    split_slide("BlueField-4 与 ICMS", [
        "#G3.5 上下文存储层",
        "介于 HBM/DRAM 与 SSD 之间",
        "每 GPU 最高 16 TB 上下文",
        "KV Cache 成为共享资源"
    ], IMG_KV_CACHE,
    """最后一个秘密武器——BlueField-4。

（形象解释）
我们来看这张图。这是一个巨大的图书馆，存放着AI的所有对话记忆。

以前，每个GPU都有自己的小书架，用完就扔。现在，BlueField-4把所有书架连成一个超级图书馆。你用过的书放回去，别人还能借。

16 TB！每个GPU可以访问16 TB的上下文存储。

还记得那个写报告写到一半被AI遗忘的同事吗？有了这个，AI可以记住你说的每一句话。从第一个字到最后一个字。""")

    # --- P16: Rubin Summary ---
    content_slide("Rubin 硬件：小结", [
        "#核心理念：物理统一内存空间",
        "22 TB/s HBM4 + 1.8 TB/s C2C",
        "每 GPU 54 TB + 16 TB 上下文",
        "#代价：全栈锁定，高资本投入"
    ], """（人设塑造：真诚的奋斗者）

让我来总结一下NVIDIA的方案。

核心就一句话：用最强的硬件，造最大的共享内存池。

22 TB/s的GPU带宽，1.8 TB/s的芯片互联，每GPU高达70 TB的可访问内存。三道防线，层层递进。

但说实话，代价也不小。首先，贵。一套NVL72，几百万美元起步。其次，锁定。一旦选择NVIDIA，就很难换别家。整个生态都绑在一起。

这是巨头的游戏。有钱，才能玩得转。

那如果你没有那么多预算，怎么办？接下来，让我们看另一条路。""")

    # --- P17: DeepSeek Engram ---
    split_slide("DeepSeek Engram", [
        "#软件定义的外挂内存",
        "无需专有硬件",
        "利用主机 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM,
    """现在，让我们换个视角。

（情景带入）
想象你是一个刚毕业的创业者。你有一个改变世界的AI梦想，但银行卡余额只有五位数。NVIDIA的方案？想都别想。

这时候，DeepSeek出现了。

他们说："朋友，我们理解你的处境。来，我们教你一招——用软件的方法，让便宜的DDR内存，变成GPU的外挂记忆。"

不需要专有硬件，普通服务器就能用。这就是Engram。

软件的力量，有时候比硬件更强大。""")

    # --- P18: Core Insight ---
    content_slide("核心洞察", [
        "LLM 大量计算浪费在重构静态知识",
        "为什么不直接查表?",
        "## 重构 vs. 查表"
    ], """DeepSeek的洞察是什么？

（生活化比喻）
想象你在背课文。老师让你复述《岳阳楼记》，你每次都要从头默背一遍。但其实你早就记住了，何必每次都重新背呢？直接翻开书，找到那一页不就行了？

大语言模型也是这样！它90%的计算，都在"重新回忆"那些早已知道的东西。"北京是中国首都"、"水的化学式是H2O"——这些每次都要重新计算，太浪费了！

为什么不直接查表呢？

这就是Engram的精髓：别算了，直接查！""")

    # --- P19: N-gram Hashing ---
    split_slide("N-gram Hashing", [
        "#输入 Token 编码为 Hash",
        "确定性、上下文相关",
        "用 Hash 作为索引查询",
        "## O(1) 复杂度"
    ], IMG_NGRAM,
    """具体怎么查？

（生活化比喻）
想象你有一本超级大字典，1000亿个词条。现在我说"北京"，你要找到关于北京的解释。

如果从头翻到尾？翻到猴年马月也翻不完。

聪明的办法：给每个词条一个编号，按编号存放。我说"北京"，你算出编号是12345，直接翻到第12345页。

这就是N-gram Hashing。把输入的词转成一个数字，用这个数字去查表。

关键是什么？O(1)复杂度！无论表有多大，查询时间都是恒定的。

一秒都不用等。""")

    # --- P20: Embedding Table ---
    data_slide("知识仓库", "100B", "参数量", "存放在主机 DDR 内存中",
    """100B参数的知识表！

（比喻）
把大语言模型想象成一个学生。

传统的大模型是天才学霸，脑子特别大（参数多），什么都记在脑子里。

Engram的模型呢，是个聪明的普通学生。脑子一般大（27B参数），但他带了一本超厚的参考书（100B知识表）。考试的时候，脑子里想不起来的，翻书就行。

结果呢？聪明学生考得和学霸一样好，甚至更好！因为书上的东西比脑子记得更准确。

27B的小模型 + 100B的知识表，效果堪比万亿参数大模型。""")

    # --- P21: Prefetching ---
    split_slide("确定性预取", [
        "#PCIe 带宽远低于 NVLink",
        "解决方案：提前获取",
        "计算 Hash → 异步加载",
        "## 吞吐损失 < 3%"
    ], IMG_PREFETCH,
    """但有个问题：从主机DDR取数据，比GPU内部慢得多。

（比喻）
就像你考试时，参考书在隔壁教室。跑去拿书要5分钟，等你回来考试都结束了。

DeepSeek怎么办？提前派人去拿！

因为需要查哪些知识是可以提前算出来的。所以在考试开始之前，就派人去图书馆借书。等你真正需要的时候，书已经送到手边了。

这张图展示的就是这个过程：GPU在算第一步的时候，内存已经在准备第三步的数据了。完美重叠。

结果呢？吞吐量损失不到3%！便宜的DDR内存，变成了高速HBM的平替。""")

    # --- P22: Engram vs RAG ---
    content_slide("Engram vs RAG", [
        "#传统 RAG：文档级检索",
        "在模型外部操作",
        "#Engram：Token 级查表",
        "在模型内部融合",
        "粒度更细，延迟更低"
    ], """有人问：这和RAG有什么区别？

（对比）
传统RAG像查百科全书。你问"北京的天气"，它把整篇"北京"的词条都给你，然后你自己找答案。

Engram像问一个无所不知的助手。你问"北京今天温度"，它直接告诉你"23度"。精确到你需要的那个点。

粒度完全不同。

而且RAG是在模型外面操作的，先检索再推理，两步走。Engram嵌入模型内部，检索和推理同时进行。

快，而且准。""")

    # --- P23: Engram Summary ---
    content_slide("Engram 软件：小结", [
        "#核心理念：算法替代硬件",
        "N-gram Hashing + 确定性预取",
        "#代价：需修改模型架构",
        "低成本，高可及性"
    ], """（人设塑造）

总结一下DeepSeek的方案。

核心理念：用算法的智慧，弥补硬件的不足。

N-gram Hashing让查询飞快，确定性预取让延迟隐藏。两个简单的技巧，组合出惊人的效果。

代价是什么？需要改模型架构。不是拿来就能用，要动手改代码。

但换来的是：低成本，高可及性。普通人也能玩得起AI了。

（真诚的脆弱性）
说实话，刚开始看到这个方案的时候，我是持怀疑态度的。真的能行吗？后来跑了实验，我才相信：软件的力量，真的可以弥补硬件的差距。""")

    # --- P24-30: Comparison Section ---
    split_slide("路线对比", [
        "#NVIDIA: Cohesion (内聚)",
        "硬件透明，模型无感知",
        "#DeepSeek: Retrieval (检索)",
        "软件显式，模型需适配"
    ], IMG_COMPARE,
    """（雷氏对比法：二元对立）

现在，让我们把两位选手放到擂台上PK。

NVIDIA的方案叫"内聚"。就像住智能公寓，你不用知道水从哪来、电从哪来，开关一开就有。

DeepSeek的方案叫"检索"。像住老房子，热水得自己烧，饭得自己做。但你知道每样东西怎么来的，可以自己改进。

一个是傻瓜模式，一个是DIY模式。

简单说：NVIDIA省心，DeepSeek省钱。""")

    content_slide("目标记忆对比", [
        "#Rubin: 动态 KV Cache",
        "对话上下文（短期记忆）",
        "#Engram: 静态 Knowledge",
        "世界知识（长期记忆）"
    ], """更关键的是，两者解决的问题不一样。

（比喻）
你的记忆分两种：一种是"刚才发生了什么"——今天早餐吃了啥、同事刚说了啥。这是短期记忆。另一种是"世界是怎样的"——中国首都是北京、1+1=2。这是长期记忆。

NVIDIA的Rubin主要解决短期记忆——让AI记住你们的对话。
DeepSeek的Engram主要解决长期记忆——让AI知道世界的事实。

所以，两者其实是互补的！最理想的方案是：用Engram存知识，用Rubin存对话。""")

    content_slide("技术参数 PK", [
        "#NVIDIA Rubin",
        "22 TB/s 带宽, 54 TB 内存, 专有硬件",
        "#DeepSeek Engram",
        "PCIe 带宽, 无限 DDR, 开源软件"
    ], """（以己之长攻彼之短）

从参数看：Rubin有22 TB/s带宽、54 TB内存，性能碾压。但需要专有硬件，起步价几百万。

Engram只有PCIe带宽，但DDR内存想加多少加多少，成本极低。

这就像比较高铁和自行车。高铁快？当然快。但自行车也能到目的地，而且不用买票。

你选哪个，取决于你的需求和预算。""")

    content_slide("适用场景", [
        "#Rubin 适合",
        "超大规模云厂商 (AWS/Azure)",
        "#Engram 适合",
        "资源受限，追求性价比"
    ], """谁该选哪个？

如果你是微软、亚马逊这种大厂，财大气粗，追求极致性能和SLA保障——选Rubin，没毛病。

如果你是创业公司、研究机构，或者像我这样的个人开发者，预算有限但梦想无限——选Engram。

没有最好的，只有最合适的。""")

    content_slide("互补而非替代", [
        "Rubin 解决 KV Cache",
        "Engram 解决知识存储",
        "## 两者可以结合"
    ], """我想强调一个观点：这两个方案不是非此即彼的关系，而是可以结合的！

最聪明的做法是什么？两手抓！

用Engram把静态知识卸载到DDR，给HBM腾出空间。用Rubin的技术管理动态KV Cache。

DDR存知识，HBM存对话。各取所长。

这，才是最优解。""")

    content_slide("开源 vs 闭源", [
        "#NVIDIA",
        "专有生态，高壁垒，高锁定",
        "#DeepSeek",
        "开源权重，低门槛，普惠创新"
    ], """从生态看：NVIDIA是封闭花园，围墙高筑。进去容易，出来难。

DeepSeek是开放草原，人人可以来，人人可以建设。

一边是商业壁垒，一边是开源共享。

这不仅是技术之争，更是理念之争。你愿意为便利付出自由，还是愿意为自由付出便利？""")

    # --- P31-34: Performance & Economics ---
    data_slide("Engram 性能", "+5.0", "BBH 推理测试", "27B 模型超越同等 MoE",
    """让数据说话。

在BBH复杂推理测试中，Engram比同规模MoE模型高出5分！

这是什么概念？就像高考多考5分，可能从二本变一本。

一个27B的小模型，打败了那些160B的巨无霸。

这就是算法的力量——以小博大，以弱胜强。""")

    data_slide("Rubin 经济效益", "10x", "Token 成本下降", "vs Blackwell，投资回报 50:1",
    """再看Rubin的经济账。

Token成本下降10倍！

NVIDIA估算：每投入1亿美元基础设施，能产出50亿美元收入。50倍的ROI！

这就是为什么大厂抢着买。对他们来说，这不是成本，这是印钞机。""")

    content_slide("AI 民主化", [
        "无需 H100/Rubin",
        "RTX 4090 集群跑 GPT-5",
        "## 打破 NVIDIA 税"
    ], """但对普通人呢？

Engram带来了AI的民主化。

想象一下：你不需要排队等H100，不需要付高昂的云服务费。用几块RTX 4090组个集群，就能跑出GPT-5级别的效果。

这就像AI界的"打土豪分田地"——打破NVIDIA税，让每个人都能参与AI革命。

这，才是真正的普惠。""")

    content_slide("地缘政治", [
        "面对出口管制",
        "Engram 是突围方案",
        "DeepSeek V3 仅用 $5.58M",
        "## More with Less"
    ], """还有一层更深的意义。

在芯片出口管制的背景下，很多地方拿不到最新GPU。Engram提供了突围之路——用软件弥补硬件的差距。

DeepSeek V3仅用558万美元完成训练，而GPT-4据说花了1亿。

这是"以少胜多"的智慧。环境越艰苦，创新越强劲。""")

    # --- P35-38: Future Section ---
    content_slide("未来：Agentic AI", [
        "从 Chatbot 到 Agent",
        "多步推理，长期记忆",
        "管理复杂工作流",
        "## 上下文是新货币"
    ], """最后，让我们看向未来。

AI正在从"聊天机器人"进化成"智能代理"。未来的AI不只回答问题，要帮你完成复杂任务：订机票、做PPT、写代码、管项目。

这需要什么？长期记忆，理解上下文。

上下文，将成为AI时代的新货币。谁的上下文长，谁的AI就更强。""")

    content_slide("百万 Token 上下文", [
        "Rubin CPX 摄入整个代码库",
        "一小时视频作为输入",
        "## 仓库即记忆"
    ], """想象这个场景：

你让AI帮你重构代码。它把整个代码仓库——几十万行——全部读入记忆，理解每个函数、每个类，然后给你完美的重构方案。

或者，你让它剪辑一小时的会议录像，它记住每一分钟的内容，自动找出重点。

这就是百万Token上下文的威力。""")

    content_slide("共享记忆池", [
        "BlueField-4 管理 KV Cache",
        "多 Agent 共享，无需重算",
        "## 协作式 AI"
    ], """更进一步，多个AI可以共享同一份记忆。

一个AI生成的上下文，另一个AI直接拿来用，不用重新计算。

就像一个团队共享同一个知识库。

这意味着：AI之间可以协作了！一个负责搜索，一个负责分析，一个负责写作。协作式AI的时代，即将到来。""")

    content_slide("长期一致性", [
        "Engram 三种记忆类型",
        "Episodic / Semantic / Procedural",
        "## +15 分，仅用 1% Token"
    ], """Engram还支持三种记忆类型：情景记忆、语义记忆、程序记忆。就像人脑一样。

在记忆一致性测试中，比基线高15分，而且只用了1%的Token！

少说话，多记事。这就是高效记忆的威力。""")

    # --- P39: Final Synthesis ---
    content_slide("终极合成", [
        "#硬件派：物理统一内存",
        "专有硅片，极致性能",
        "#软件派：算法定义内存",
        "数学技巧，极致效率"
    ], """在结束之前，让我做一个终极对比。

NVIDIA代表的是"集成电路"思维——用专有硅片创造物理统一的内存空间。22 TB/s带宽，1.8 TB/s互联，这是硬件能力的极限。代价是全栈锁定和高额投入。

DeepSeek代表的是"算法"思维——用数学技巧把普通DDR变成高效的知识存储。N-gram Hashing、确定性预取，这是软件智慧的结晶。代价是需要改架构。

两条路，两种哲学。一个是"钞能力"，一个是"脑能力"。

但无论哪条路，目标只有一个：让AI不再遗忘。""")

    # --- P40: Summary (峰终定律：高潮) ---
    content_slide("总结：内存革命", [
        "Rubin 定义了性能天花板",
        "Engram 拉高了可及地板",
        "## 殊途同归：治愈 AI 失忆"
    ], """（峰终定律：高潮总结）

朋友们，让我来做个总结。

今天我们讲了一场战争——内存的战争。

NVIDIA用钞能力定义了性能的天花板。有钱能使鬼推磨，有钱也能让AI永不遗忘。

DeepSeek用智慧拉高了可及性的地板。没钱也能玩AI，普通人也能参与革命。

两条不同的路，殊途同归。都在治愈同一个病——AI的失忆症。

2026年，将被铭记为"内存革命元年"。从今天起，我们告别内存贫困的时代，迎来记忆永驻的新纪元。

这不是我一个人的观点，这是整个行业的共识。""")

    # --- P40: Closing (峰终定律：结尾要情怀) ---
    closing_slide("谢谢", "The Architecture of Memory",
    """（峰终定律：情怀收尾）

朋友们，这就是今天分享的全部内容。

故事讲完了，但历史才刚刚开始。

无论你选择硬件的霸道，还是软件的智慧；无论你是财大气粗的巨头，还是白手起家的创业者——这都是一个属于我们每个人的新时代。

AI不再会遗忘。它将永远记得你。

（悬念收尾）
这只是开始。明年，我们还会有更大的惊喜。

感谢大家的聆听！朋友们，我们下次再见！""")

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_V9_leijunskill2.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_ppt_v9()
