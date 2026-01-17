"""
Memory Architecture Report - V8 (40 Pages with Storytelling Speaker Notes)
Generated following leijunskill standards:
- 40 pages with storytelling-enhanced Lei Jun style speaker notes
- Each note includes: life metaphors, scenario immersion, number visualization, contrast stories
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- IMAGE PATHS ---
IMG_TITLE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_cover_memory_wall_1768527608448.png"
IMG_RUBIN = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_rubin_chip_1768527625522.png"
IMG_ENGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_engram_software_1768527643236.png"
IMG_COMPARE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/ppt_comparison_split_1768527658633.png"
IMG_MEMORY_WALL = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_memory_wall_1768530473358.png"
IMG_KV_CACHE = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_kv_cache_1768530494022.png"
IMG_NGRAM = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_ngram_hashing_1768530510183.png"
IMG_PREFETCH = r"C:/Users/mgcpn/.gemini/antigravity/brain/362d749a-8462-4250-84e0-af5da822d504/concept_prefetch_strategy_1768530526058.png"

# --- CONSTANTS ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BG_COLOR = RGBColor(15, 15, 25)
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(1.5)
CONTENT_TOP = Inches(2.0)
CONTENT_HEIGHT = Inches(4.0)
SPLIT_TEXT_LEFT = Inches(0.5)
SPLIT_TEXT_WIDTH = Inches(5.5)
SPLIT_IMG_LEFT = Inches(6.5)
SPLIT_IMG_WIDTH = Inches(6.3)


def create_ppt_v8():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes_text

    def add_signature(slide, is_cover=True):
        sig_width = Inches(5)
        sig_height = Inches(0.9)
        sig_left = (SLIDE_WIDTH - sig_width) / 2
        sig_top = SLIDE_HEIGHT - sig_height - Inches(0.15)
        box = slide.shapes.add_textbox(sig_left, sig_top, sig_width, sig_height)
        tf = box.text_frame
        p1 = tf.add_paragraph()
        p1.text = "天翼云湖北分公司" + (" 作者" if is_cover else "")
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(0, 191, 255)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "王理"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(255, 215, 0)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def add_title(slide, text, center=False):
        box = slide.shapes.add_textbox(Inches(0.5), TITLE_TOP, SLIDE_WIDTH - Inches(1), TITLE_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        if center:
            p.alignment = PP_ALIGN.CENTER

    def add_content(slide, lines):
        box = slide.shapes.add_textbox(Inches(1), CONTENT_TOP, SLIDE_WIDTH - Inches(2), CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(16)
            if line.startswith("##"):
                p.text = line[2:].strip()
                p.font.size = Pt(54)
                p.font.color.rgb = RGBColor(255, 215, 0)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            elif line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = "• " + line
                p.font.size = Pt(26)
                p.font.color.rgb = RGBColor(220, 220, 220)

    def add_split(slide, lines, img_path):
        box = slide.shapes.add_textbox(SPLIT_TEXT_LEFT, CONTENT_TOP, SPLIT_TEXT_WIDTH, CONTENT_HEIGHT)
        tf = box.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.space_after = Pt(12)
            if line.startswith("#"):
                p.text = line[1:].strip()
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
            else:
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(220, 220, 220)
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, SPLIT_IMG_LEFT, CONTENT_TOP, width=SPLIT_IMG_WIDTH)

    # ========== PAGE GENERATION (40 pages with storytelling notes) ==========

    # --- P1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "内存的战争\nThe Architecture of Memory", center=True)
    if os.path.exists(IMG_TITLE):
        s.shapes.add_picture(IMG_TITLE, Inches(4), Inches(2.2), height=Inches(3.0))
    sub = s.shapes.add_textbox(Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "NVIDIA Rubin vs. DeepSeek Engram | 2026 内存架构深度解析"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=True)
    add_notes(s, """朋友们，大家好！

今天，我想和大家讲一个故事。这个故事的主角，是"记忆"。

想象一下：你正在和一个AI助手聊天。你们聊了半个小时，从人生理想聊到今晚吃什么。突然，你问它："我刚才说的那个餐厅叫什么来着？" 它愣住了，一脸茫然地回答："抱歉，你说过餐厅吗？"

这就是AI的失忆症。今天，我们要聊的就是：如何治好它。

两个完全不同的阵营，给出了两个截然相反的答案。这，就是内存的战争。""")

    # --- P2: Agenda ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "汇报大纲")
    add_content(s, [
        "1. 时代的挑战：两堵高墙",
        "2. NVIDIA Rubin：硬件霸权",
        "3. DeepSeek Engram：软件革命",
        "4. 路线对比：硬件 vs 软件",
        "5. 性能与经济分析",
        "6. 未来展望：Agentic AI"
    ])
    add_notes(s, """我们今天的分享，就像一场探险。

首先，我们会遇到两堵高墙——内存墙和上下文墙。这是AI发展路上最大的拦路虎。

然后，我们会看到两个英雄如何翻越这堵墙。一个是财大气粗的硬件巨头NVIDIA，另一个是聪明绝顶的软件新秀DeepSeek。

接着，我们会站在山顶，对比两条路线的风景。

最后，我们一起眺望远方，看看AI的未来会是什么模样。

准备好了吗？Let's go！""")

    # --- P3: Paradigm Shift ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2026：范式转移")
    add_content(s, [
        "从算力 (FLOPS) 到 内存 (Memory)",
        "万亿参数模型成为常态",
        "## 记忆，是新的石油"
    ])
    add_notes(s, """朋友们，你们有没有注意到一个变化？

过去这些年，我们一直在追逐算力。谁的GPU更快，谁的FLOPS更高，谁就是王者。

但现在不一样了。就像当年石油取代煤炭成为工业血液一样，今天，记忆正在取代算力，成为AI的新命脉。

为什么？因为模型参数已经突破了万亿大关。你可以把算力想象成一个超级大厨，厨艺精湛，切菜飞快。但问题是——冰箱太小了！食材放不下，大厨再厉害也只能干瞪眼。

记忆，就是那个冰箱。谁的冰箱大，谁就能做更多菜。""")

    # --- P4: Memory Wall ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Memory Wall：内存墙")
    add_split(s, [
        "#GPU 算力远超内存带宽",
        "数据传输成为瓶颈",
        "权重访问速度限制推理",
        "#算力越强，内存墙越高"
    ], IMG_MEMORY_WALL)
    add_notes(s, """我们来看这张图。

想象一下早高峰的地铁站。左边是源源不断涌来的乘客——这就是GPU产生的算力。右边是站台——这就是内存里的数据。中间这道闸机，就是带宽。

问题来了：乘客太多，闸机太少。大家挤在那里，谁也过不去。

这就是内存墙。GPU算力再强，数据传不过来，也只能空转。更可怕的是：算力越强，乘客越多，闸机还是那几个，堵得更厉害。

这是一个让工程师们抓狂的困境。""")

    # --- P5: Context Wall ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Context Wall：上下文墙")
    add_content(s, [
        "KV Cache 随序列长度指数膨胀",
        "模型无法记住完整对话",
        "处理百万 Token 需天价计算",
        "## 失忆，是 AI 最大的痛点"
    ])
    add_notes(s, """第二堵墙更扎心——上下文墙。

我给大家讲个真实的故事。我有个朋友，用AI写小说。写到第三章的时候，AI突然把主角的名字给忘了！从"张三"变成了"李四"。他气得差点把电脑砸了。

这就是上下文墙的问题。AI的记忆是有限的。你和它聊得越多，它需要记住的东西就越多。这些记忆叫做KV Cache，会随着对话长度指数级增长。

就像你的手机相册。照片越来越多，手机越来越卡。最后不得不删照片。AI也一样，聊着聊着，就把你之前说的忘了。

失忆，是AI最大的痛点。我们必须解决它。""")

    # --- P6: Two Camps ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "两大阵营的回答")
    add_content(s, [
        "#NVIDIA (硬件派)",
        "Rubin 架构：极致协同，专有互联",
        "#DeepSeek (软件派)",
        "Engram 技术：算法优化，软件定义内存"
    ])
    add_notes(s, """面对这两堵墙，世界上出现了两个阵营。

一边是NVIDIA，就像一个土豪。他们说："墙太高？造个电梯！钱不是问题，用最贵的材料，造最大的梯子。"

另一边是DeepSeek，像一个聪明的穷学生。他们说："我们没钱买电梯，但我们可以找条隧道绕过去啊！"

一个靠钱，一个靠脑子。两条完全不同的路，目标却是一样的——翻过那堵墙。

接下来，让我们分别看看他们是怎么做的。""")

    # --- P7: Rubin Architecture ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVIDIA Rubin 架构")
    add_split(s, [
        "#Jensen Huang @ CES 2026",
        "年度节奏发布",
        "Extreme Co-design (极致协同)",
        "六大芯片生态系统"
    ], IMG_RUBIN)
    add_notes(s, """让我们先看土豪的方案。

2026年CES上，老黄穿着他标志性的皮衣走上台。他说："朋友们，我们不是在做一块GPU，我们是在做一个超级大脑。"

Rubin架构的思路很简单也很暴力：既然问题是内存不够，那就把所有东西都连起来，共享内存。GPU、CPU、网卡、交换机……全部融为一体。

就像把一栋栋独立的别墅，改造成一个超大的公寓楼。每个房间都能用整栋楼的水电。

这叫"极致协同"。听起来简单，做起来要命。""")

    # --- P8: Six Chip Ecosystem ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "六大芯片生态系统")
    add_content(s, [
        "Rubin GPU - 核心计算引擎",
        "Vera CPU - 内存扩展",
        "NVLink 6 Switch - 高速互联",
        "ConnectX-9 SuperNIC - 网络加速",
        "BlueField-4 DPU - 上下文管理",
        "Spectrum-6 Ethernet Switch"
    ])
    add_notes(s, """这六块芯片，就像复仇者联盟。

Rubin GPU是钢铁侠，负责战斗力，算力最强。
Vera CPU是绿巨人，力大无穷，扛着1.5TB的内存。
NVLink 6是蜘蛛侠的蛛丝，把所有人连在一起。
ConnectX-9是鹰眼，眼观六路，网络通信全靠它。
BlueField-4是幻视，管理着所有人的记忆。
Spectrum-6是黑寡妇，在后台默默协调一切。

六个超级英雄，组成一支无敌战队。这就是NVIDIA的全明星阵容。""")

    # --- P9: Rubin GPU ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin GPU：性能怪兽")
    add_content(s, [
        "## 3360 亿",
        "晶体管数量 (台积电 3nm)",
        "## 50 PFLOPS",
        "FP4 推理性能"
    ])
    add_notes(s, """大家看这个数字——3360亿！

3360亿是什么概念？如果每个晶体管是一粒沙子，这些沙子可以填满40个奥运会游泳池！全部塞进一块比你手掌还小的芯片里。

再看这个——50 PFLOPS。

这是什么概念？2008年，全世界最快的超级计算机叫"走鹃"，峰值性能1 PFLOPS，占地几百平米，耗电几兆瓦。今天，一块Rubin GPU，装在你的电脑里，就有50个"走鹃"的算力。

17年，50倍。科技进步的速度，比我们想象的快得多。""")

    # --- P10: HBM4 Bandwidth ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：速度的革命")
    add_content(s, [
        "## 22 TB/s",
        "内存带宽 (vs Blackwell 8 TB/s)",
        "## 2.75x",
        "代际提升"
    ])
    add_notes(s, """22 TB/s！

这是什么概念？让我给大家算一笔账。

一部高清电影大概5GB。22 TB/s意味着：每秒钟可以传输4400部高清电影。一眨眼的功夫，整个Netflix的片库就传完了。

还记得刚才说的地铁闸机吗？以前是8个闸机，现在变成22个，还都是快速通道。乘客再多，也不用挤了。

这不是小改进，这是代际飞跃。NVIDIA用硬实力，硬生生把那堵墙推倒了一半。""")

    # --- P11: HBM4 Capacity ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "HBM4：容量震撼")
    add_content(s, [
        "## 288 GB",
        "单卡 HBM4 显存容量",
        "可容纳 1440 亿参数模型权重",
        "MoE 路由决策可在微秒内完成"
    ])
    add_notes(s, """288 GB！单卡288 GB！

我给大家讲个故事。三年前，我们团队要跑一个700亿参数的模型。买了8块A100，每块80GB，凑了640GB才勉强够用。机器塞满了一整个机柜，电费一个月好几万。

现在呢？一块Rubin，288GB，装下1440亿参数，一块卡就够了。

这就像以前搬家需要一整个车队，现在一辆SUV就能装下所有家当。

技术的进步，就是这么简单粗暴。""")

    # --- P12: Vera CPU ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Vera CPU：AI 工厂专属")
    add_content(s, [
        "88 颗 Olympus 定制 ARM 核心",
        "## 1.5 TB",
        "LPDDR5X 内存池",
        "首款支持 FP8 的 CPU"
    ])
    add_notes(s, """但288 GB还不够。NVIDIA又造了一个"外挂粮仓"——Vera CPU。

1.5 TB的内存，是GPU显存的5倍多。

想象一下：GPU是一个超级大厨，手边的料理台（HBM）放不下那么多食材。没关系，旁边还有一个巨大的冷库（Vera内存），什么都往里塞。需要的时候，通过高速传送带（NVLink）瞬间送过来。

关键是，这个冷库就在大厨隔壁，不是隔着几条街的超市。拿东西快得很。

88颗定制核心，就为一件事——快速搬运数据。这就是Vera CPU的使命。""")

    # --- P13: NVLink-C2C ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVLink-C2C：高速一致性")
    add_content(s, [
        "## 1.8 TB/s",
        "芯片间互联带宽",
        "相比 Grace-Blackwell 提升 2x",
        "GPU 可像访问本地内存一样访问 CPU 内存"
    ])
    add_notes(s, """1.8 TB/s的芯片间互联！

这是什么概念？从GPU到CPU的数据传输速度，和GPU内部访问自己显存的速度，几乎一样快。

就像两栋楼之间架了一座天桥，走路过去只要10秒钟。以前呢？要下楼、过马路、等红灯、上楼，至少10分钟。

对于程序来说，它甚至感觉不到内存在两个不同的芯片上。就像住在公寓里，你不需要知道水是从哪个水厂来的，打开水龙头就有水。

这就是硬件工程师的魔法——把复杂的事情，变得简单透明。""")

    # --- P14: NVL72 ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "NVL72：超级机柜")
    add_content(s, [
        "72 块 Rubin GPU + 36 块 Vera CPU",
        "## 54 TB",
        "机柜总内存容量",
        "统一内存寻址，透明访问"
    ])
    add_notes(s, """把所有这些加起来，就是NVL72超级机柜。

54 TB的内存！

这是什么概念？假设一个人的大脑记忆容量是2.5 PB（一个流行的估算），那么NVL72相当于人类大脑的2%。听起来不多？但别忘了，人脑花了几百万年进化，NVL72只用了几年开发。

更重要的是：这54 TB的内存，对于所有GPU来说都是透明可见的。就像一个超大的共享硬盘，谁需要就去取，不需要问别人。

这就是NVIDIA的暴力美学——用最贵的硬件，堆出最强的性能。""")

    # --- P15: BlueField-4 ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "BlueField-4 与 ICMS")
    add_split(s, [
        "#G3.5 上下文存储层",
        "介于 HBM/DRAM 与 SSD 之间",
        "每 GPU 最高 16 TB 上下文",
        "KV Cache 成为共享资源"
    ], IMG_KV_CACHE)
    add_notes(s, """最后一个秘密武器——BlueField-4和ICMS。

我们来看这张图。这是一个巨大的图书馆，里面存放着所有的对话记忆（KV Cache）。

以前，每个GPU都有自己的小书架，用完就扔。现在，BlueField-4把所有书架连成一个超级图书馆。你用过的书，放回去，别人还能借。

16 TB！每个GPU可以访问16 TB的上下文存储。

还记得刚才那个写小说忘记主角名字的朋友吗？有了这个，AI可以记住你写的每一个字，从第一章到第一百章，一个都不会忘。

这，就是治愈AI失忆症的终极方案。""")

    # --- P16: Rubin Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Rubin 硬件总结")
    add_content(s, [
        "#核心理念：物理统一内存空间",
        "22 TB/s HBM4 + 1.8 TB/s C2C + 16 TB ICMS",
        "#代价：全栈锁定 NVIDIA 生态",
        "高 CapEx，高性能"
    ])
    add_notes(s, """让我来总结一下NVIDIA的方案。

核心理念只有一个：造一个超级大的共享内存池，让所有芯片都能用。

22 TB/s的GPU带宽，1.8 TB/s的芯片互联，16 TB的上下文存储。三个数字，三道防线，层层递进。

代价是什么？首先，贵。一套NVL72系统，几百万美元起步。其次，锁定。一旦选择NVIDIA，就很难换别家。

这是巨头的游戏。有钱，任性。

但如果你没有那么多预算呢？别急，接下来我们看看另一条路。""")

    # --- P17: DeepSeek Engram ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "DeepSeek Engram")
    add_split(s, [
        "#软件定义的外挂内存",
        "无需专有硬件",
        "利用主机 DDR 内存",
        "打破显存容量限制"
    ], IMG_ENGRAM)
    add_notes(s, """现在，让我们换一个场景。

想象一下：你是一个刚毕业的创业者，怀揣着AI梦想，但银行卡余额只有五位数。NVIDIA的方案？想都别想。

这时候，DeepSeek出现了。他们说："朋友，我们理解你。来，我们教你一招。"

Engram的思路完全不同：既然买不起更大的显存，那就用主机上便宜的DDR内存。反正都是内存，想办法让GPU能用上就行。

就像租不起大房子，那就在旁边租个仓库，需要的时候去取东西。关键是怎么取得快。

这就是软件的智慧——用巧劲，四两拨千斤。""")

    # --- P18: Core Insight ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "核心洞察")
    add_content(s, [
        "LLM 大量计算浪费在重构静态知识",
        "为什么不直接查找?",
        "## Reconstruct vs. Lookup",
        "用查表替代计算"
    ])
    add_notes(s, """DeepSeek的核心洞察是什么？他们发现了一个惊人的事实。

想象一下你在背课文。老师让你复述《岳阳楼记》，你每次都要从头默背一遍。但其实，你早就记住了，何必每次都重新背呢？直接翻开书，找到那一页不就行了？

大语言模型也是这样。它90%的计算，都在"重新回忆"那些早就知道的东西。"北京是中国的首都"、"水的化学式是H2O"——这些知识每次都要重新计算一遍，太浪费了！

为什么不直接查表呢？把这些知识存在一个大表里，需要的时候查一下，多快啊！

这就是Engram的精髓：别算了，直接查！""")

    # --- P19: N-gram Hashing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "N-gram Hashing 原理")
    add_split(s, [
        "#输入 Token 序列编码为 Hash",
        "确定性、上下文相关的 ID",
        "用 Hash 作为 Key 查询",
        "## O(1) 复杂度查表"
    ], IMG_NGRAM)
    add_notes(s, """具体怎么查呢？我们来看这张图。

想象你有一本超级大的字典，里面有1000亿个词条。现在我说"北京"，你要在字典里找到关于北京的解释。

如果从头翻到尾？那得翻到猴年马月。

聪明的办法是：给每个词条一个编号，按编号存放。我说"北京"，你算出编号是12345，直接翻到第12345页，一秒搞定。

这就是N-gram Hashing。把输入的词转成一个数字（Hash），用这个数字去查表。无论表有多大，查询时间都是恒定的——O(1)复杂度。

快得让人难以置信。""")

    # --- P20: Embedding Table ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Embedding Table：知识仓库")
    add_content(s, [
        "## 100B 参数",
        "存放在主机 DDR 内存中",
        "27B 活跃模型 + 100B 知识表",
        "小模型拥有万亿知识"
    ])
    add_notes(s, """100B参数的知识表！

我给大家讲个形象的比喻。

把大语言模型想象成一个学生。传统的大模型，就像一个天才学霸，什么都记在脑子里，考试的时候全靠回忆。脑子越大（参数越多），记得越多。

Engram的模型呢，像一个聪明的普通学生。他脑子一般大（27B参数），但他带了一本超厚的参考书（100B知识表）。考试的时候，脑子里想不起来的，翻书就行。

结果呢？聪明学生考得和天才学霸一样好，甚至更好。因为书里的东西比脑子记得更准确。

这就是Engram的魔法：小脑袋，大智慧。""")

    # --- P21: PCIe Bottleneck ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "克服 PCIe 瓶颈")
    add_split(s, [
        "#PCIe 带宽远低于 NVLink",
        "解决方案：确定性预取",
        "提前计算 Hash ID",
        "异步获取数据"
    ], IMG_PREFETCH)
    add_notes(s, """但有一个问题：主机内存离GPU太远了，取数据慢。

这就像你家旁边有个图书馆，但图书馆在山那边。走过去要半小时，等你把书拿回来，考试早结束了。

DeepSeek怎么解决的？提前预约！

他们发现，需要查哪些知识，是可以提前算出来的。所以在考试开始之前，就派人去图书馆借书。等你真正需要的时候，书已经送到手边了。

这张图展示的就是这个过程：上面是GPU在计算，下面是内存在预取。两条线完美同步，数据永远提前准备好。

聪明吧？这就是算法的力量。""")

    # --- P22: Prefetch Effect ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "预取策略效果")
    add_content(s, [
        "在 GPU 执行早期层时",
        "异步从 DDR 获取后续所需数据",
        "## 吞吐量损失 < 3%",
        "廉价 DDR 成为 HBM 的可行扩展"
    ])
    add_notes(s, """效果如何呢？

吞吐量损失不到3%！

让我给大家翻译一下这个数字的意义。

NVIDIA用几百万美元的专有硬件，实现了高速内存访问。DeepSeek用几百块钱的DDR内存条，加上一套巧妙的软件，达到了97%的效果。

就像一个是开法拉利，一个是骑改装电动车。法拉利快？当然快。但电动车也能到达目的地，而且省了99%的油钱。

对于大多数人来说，这97%的效果，足够了。而且，DDR内存你想加多少就加多少，便宜得很。

这就是普惠创新的力量。""")

    # --- P23: Engram vs RAG ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram vs RAG")
    add_content(s, [
        "#传统 RAG：文档级检索",
        "在 Transformer 外部操作",
        "#Engram：Token 级查表",
        "在 Transformer 内部融合",
        "粒度更细，延迟更低"
    ])
    add_notes(s, """有人会问：这和RAG有什么区别？

让我打个比方。

传统RAG就像查百科全书。你问"北京的天气"，它把整篇"北京"的词条都给你，然后你自己在里面找答案。信息太多，很多是没用的。

Engram呢，像问一个无所不知的助手。你问"北京今天温度"，它直接告诉你"23度"。精确到你需要的那个信息点，不多不少。

粒度不同，效率天差地别。

而且，RAG是在模型外面操作的，要先检索再推理，两步走。Engram是嵌入到模型内部的，检索和推理同时进行，更快。""")

    # --- P24: Engram Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Engram 软件总结")
    add_content(s, [
        "#核心理念：算法替代硬件",
        "N-gram Hashing + 预取重叠",
        "#代价：需修改模型架构",
        "低成本，高可及性"
    ])
    add_notes(s, """总结一下DeepSeek的方案。

核心理念：用算法的智慧，弥补硬件的不足。N-gram Hashing让查询飞快，预取策略让延迟隐藏。

代价是什么？需要修改模型架构。不是拿来就能用的，要动手改代码。

但好处巨大：低成本，高可及性。普通人也能玩得起。

这是一条属于普通人的AI之路。

好了，两位选手都介绍完了。接下来，让我们把他们放在一起，好好比一比。""")

    # --- P25-40: 对比、经济、未来、总结部分（同样加上故事性解说词）---
    
    # --- P25: Route Comparison ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "路线对比")
    add_split(s, [
        "#NVIDIA: Cohesion (内聚)",
        "硬件透明，模型无感知",
        "#DeepSeek: Retrieval (检索)",
        "软件显式，模型需适配"
    ], IMG_COMPARE)
    add_notes(s, """现在，让我们把两位选手放在擂台上。

NVIDIA的方案叫"内聚"。什么意思？就像住在一栋智能公寓里，你不知道水从哪里来、电从哪里来，打开开关就有。硬件在幕后默默搞定一切。

DeepSeek的方案叫"检索"。就像你住在老房子里，想用热水得自己烧，想吃饭得自己做。但是，你知道每一样东西是怎么来的，可以自己改进。

一个是傻瓜式的，一个是DIY的。各有优劣。

简单来说：NVIDIA省心，DeepSeek省钱。""")

    # --- P26: Memory Type Comparison ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "目标记忆类型对比")
    add_content(s, [
        "#Rubin: 动态 KV Cache",
        "对话上下文 (短期记忆)",
        "#Engram: 静态 Knowledge",
        "世界知识 (长期记忆)"
    ])
    add_notes(s, """更重要的是，这两个方案解决的问题不一样。

打个比方：你的记忆分两种。一种是"刚才发生了什么"——比如今天早餐吃了什么、刚才同事跟你说了什么。这是短期记忆。另一种是"世界是怎样的"——比如中国的首都是北京、1+1=2。这是长期记忆。

NVIDIA的Rubin，主要解决短期记忆问题——让AI记住你们聊了什么。
DeepSeek的Engram，主要解决长期记忆问题——让AI知道世界上的事实。

两者其实是互补的！最理想的方案是：用Engram存知识，用Rubin存对话。各司其职。""")

    # --- P27-40: 继续其他页面（简化处理）---
    
    pages_data = [
        ("技术参数对比", ["#NVIDIA Rubin", "22 TB/s 带宽, 54 TB 容量", "#DeepSeek Engram", "PCIe 带宽, 无限 DDR"],
         "从数据上看：Rubin有22 TB/s的带宽，54 TB的容量，性能碾压。但需要专有硬件，起步价几百万。Engram只有PCIe带宽，但DDR内存想加多少加多少，便宜得很。这就像比较高铁和自行车。高铁快？当然快。但自行车也能到，而且不用买票。"),
        
        ("适用场景对比", ["#Rubin 适合", "超大规模云厂商", "#Engram 适合", "资源受限环境"],
         "谁该选哪个？如果你是微软、亚马逊这样的大厂，财大气粗，追求极致性能和SLA保障，选Rubin。如果你是创业公司、研究机构、或者像我这样的个人开发者，预算有限但梦想无限，选Engram。没有最好的，只有最合适的。"),
        
        ("互补而非替代", ["Rubin 解决 KV Cache 容量", "Engram 解决知识存储", "## 两者可以结合"],
         "我想强调一点：这两个方案不是非此即彼的关系。它们解决的是不同层面的问题。最聪明的做法是什么？两手都要抓！用Engram把静态知识卸载到DDR，给HBM腾出空间。用Rubin的技术管理动态的KV Cache。各取所长，这才是最优解。"),
        
        ("开源 vs 闭源", ["#NVIDIA", "专有生态，高锁定", "#DeepSeek", "开源权重，普惠创新"],
         "从生态角度看：NVIDIA是一个封闭花园，围墙高筑。一旦进去，想出来很难。DeepSeek是一片开放草原，人人都可以来，人人都可以建设。一边是商业壁垒，一边是开源共享。这不仅是技术之争，更是理念之争。"),
        
        ("Engram 性能实测", ["## +5.0", "BBH (复杂推理)", "## +3.4", "MMLU (知识问答)"],
         "让数据说话。在BBH复杂推理测试中，Engram比同等规模的模型高出5分！这是什么概念？就像高考多考了5分，可能就从二本变成一本了。在MMLU知识问答中高出3.4分。一个27B的小模型，打败了那些160B的巨无霸。这就是算法的力量——以小博大，以弱胜强。"),
        
        ("Rubin 经济效益", ["## 10x", "Token 成本下降", "## $5B/$100M", "投资回报率"],
         "再看Rubin的经济账。Token成本下降10倍！每投入1亿美元基础设施，能产出50亿美元的收入。50倍的ROI！这就是为什么大厂们抢着买。对他们来说，这不是支出，这是印钞机。"),
        
        ("AI 民主化", ["无需 H100/Rubin", "RTX 4090 集群跑 GPT-5", "## 打破 NVIDIA 税"],
         "但对我们普通人呢？Engram带来了AI的民主化。想象一下：你不需要排队等H100，不需要支付高昂的云服务费用。用几块RTX 4090显卡组个集群，就能跑出GPT-5级别的效果。这就像是AI界的'打土豪分田地'——打破NVIDIA税，让每个人都能参与AI革命。"),
        
        ("地缘政治意义", ["面对出口管制", "Engram 是突围方案", "## More with Less"],
         "这里还有一层更深的意义。在芯片出口管制的背景下，很多地方拿不到最新的GPU。Engram提供了一条突围之路：用软件弥补硬件的差距。DeepSeek V3仅用558万美元就完成了训练，而GPT-4据说花了1亿美元。这就是'以少胜多'的智慧。"),
        
        ("未来：Agentic AI", ["从 Chatbot 到 Agent", "多步推理，长期记忆", "## Context is the new Currency"],
         "最后，让我们展望未来。AI正在从'聊天机器人'进化成'智能代理'。未来的AI不只是回答问题，而是要帮你完成复杂任务：订机票、做PPT、写代码、管理项目。这需要什么？需要长期记忆，需要理解上下文。上下文，将成为AI时代的新货币。"),
        
        ("百万 Token 上下文", ["Rubin CPX 可摄入完整代码库", "一小时视频作为输入", "## 整个仓库成为工作记忆"],
         "想象一下这个场景：你让AI帮你重构一个大型项目的代码。它把整个代码仓库——几十万行代码——全部读入记忆，理解每一个函数、每一个类的作用，然后给你一个完美的重构方案。或者，你让它帮你剪辑一个小时的会议录像，它记住每一分钟的内容，自动找出重点。这就是百万Token上下文的威力。"),
        
        ("共享记忆池", ["BlueField-4 管理 KV Cache", "多 Agent 共享，无需重算", "## 协作式 AI"],
         "更进一步，多个AI可以共享同一份记忆。一个AI生成的上下文，另一个AI直接拿来用，不用重新计算。就像一个团队共享同一个知识库。这意味着什么？AI之间可以协作了！一个负责搜索，一个负责分析，一个负责写作，无缝配合。这就是协作式AI的未来。"),
        
        ("长期一致性记忆", ["Engram 三种记忆", "Episodic/Semantic/Procedural", "## +15 分 (仅用 1% Token)"],
         "Engram还支持三种类型的记忆，就像人脑一样：情景记忆（发生了什么）、语义记忆（知道什么）、程序记忆（会做什么）。在记忆一致性测试中，比基线高出15分，而且只用了1%的Token！这就是高效记忆的威力——少说话，多记事。"),
    ]
    
    for title, content, notes in pages_data:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        add_title(s, title)
        add_content(s, content)
        add_notes(s, notes)

    # --- P39: Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "总结：内存革命")
    add_content(s, [
        "Rubin 定义了性能上限",
        "Engram 拉高了可及下限",
        "## 殊途同归：解决 AI 失忆"
    ])
    add_notes(s, """朋友们，让我来做个总结。

今天我们讲了一场战争——内存的战争。这场战争有两个主角：NVIDIA的Rubin和DeepSeek的Engram。

Rubin用钞能力定义了性能的天花板。有钱能使鬼推磨，有钱也能让AI永不遗忘。

Engram用智慧拉高了可及性的地板。没钱也能玩AI，普通人也能参与这场革命。

两条不同的路，殊途同归。它们都在解决同一个问题——治愈AI的失忆症。

2026年，将被铭记为"内存革命元年"。从今天起，我们告别了内存贫困的时代，迎来了记忆永驻的新纪元。""")

    # --- P40: Closing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    box = s.shapes.add_textbox(Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5))
    p = box.text_frame.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(0.8))
    p = sub.text_frame.add_paragraph()
    p.text = "The Architecture of Memory"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    add_signature(s, is_cover=False)
    add_notes(s, """朋友们，这就是今天分享的全部内容。

故事讲完了，但历史才刚刚开始。

无论你选择硬件的霸道，还是软件的智慧；无论你是财大气粗的巨头，还是白手起家的创业者——这都是一个属于我们每个人的新时代。

AI不再会遗忘，AI将永远记得你。

感谢大家的聆听！如果今天的分享对你有所启发，欢迎交流讨论。

朋友们，再见！""")

    # --- Save ---
    out_path = os.path.join(os.getcwd(), 'Memory_Architecture_Report_V8_storytelling.pptx')
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_ppt_v8()
